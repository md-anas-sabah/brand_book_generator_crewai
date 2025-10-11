from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.enum.text import MSO_ANCHOR
import os
import re
import sys
sys.path.append(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
from agents.font_research_agent import FontResearchAgent
from agents.introduction_content_agent import IntroductionContentAgent
from agents.brand_purpose_content_agent import BrandPurposeContentAgent
from agents.brand_story_content_agent import BrandStoryContentAgent
from agents.enhanced_color_research_agent import EnhancedColorResearchAgent
from agents.iconography_agent import IconographyAgent
from tools.fal_image_tool import generate_brand_illustrations


class GridLayout:
    """Grid-based layout system for consistent positioning"""
    
    def __init__(self, slide_width=Inches(10), slide_height=Inches(7.5), 
                 cols=12, rows=8, margin=Inches(0.5)):
        self.width = slide_width
        self.height = slide_height
        self.cols = cols
        self.rows = rows
        self.margin = margin
        
        # Calculate grid dimensions
        self.col_width = (slide_width - 2 * margin) / cols
        self.row_height = (slide_height - 2 * margin) / rows
    
    def get_position(self, col, row, width_cols=1, height_rows=1):
        """Get position and size for grid cell"""
        left = self.margin + (col * self.col_width)
        top = self.margin + (row * self.row_height)
        width = width_cols * self.col_width
        height = height_rows * self.row_height
        return left, top, width, height


class TextStyler:
    """Reusable text styling system"""
    
    def __init__(self, primary_font="Inter", secondary_font="Source Sans Pro", 
                 primary_color="#2E86AB", text_color="#FFFFFF"):
        self.primary_font = primary_font
        self.secondary_font = secondary_font
        self.primary_color = primary_color
        self.text_color = text_color
    
    def apply_title_style(self, paragraph, color=None, size=32):
        """Large, bold titles using primary font with color support"""
        paragraph.font.name = self.primary_font
        paragraph.font.size = Pt(size)
        paragraph.font.bold = True
        
        if color == 'white':
            paragraph.font.color.rgb = RGBColor(255, 255, 255)
        elif color == 'black':
            paragraph.font.color.rgb = RGBColor(0, 0, 0)
        elif isinstance(color, str) and color.startswith('#'):
            # Handle hex colors
            hex_color = color.lstrip('#')
            r, g, b = tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
            paragraph.font.color.rgb = RGBColor(r, g, b)
        else:
            paragraph.font.color.rgb = RGBColor(*self._hex_to_rgb(color or self.text_color))
    
    def apply_subtitle_style(self, paragraph, color=None, size=24):
        """Medium subtitles using secondary font with color support"""
        paragraph.font.name = self.secondary_font
        paragraph.font.size = Pt(size)
        paragraph.font.bold = False
        
        if color == 'white':
            paragraph.font.color.rgb = RGBColor(255, 255, 255)
        elif color == 'black':
            paragraph.font.color.rgb = RGBColor(0, 0, 0)
        elif isinstance(color, str) and color.startswith('#'):
            # Handle hex colors
            hex_color = color.lstrip('#')
            r, g, b = tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
            paragraph.font.color.rgb = RGBColor(r, g, b)
        else:
            paragraph.font.color.rgb = RGBColor(*self._hex_to_rgb(color or self.text_color))
    
    def apply_body_style(self, paragraph, color=None, size=16):
        """Body text using secondary font with color support"""
        paragraph.font.name = self.secondary_font
        paragraph.font.size = Pt(size)
        paragraph.font.bold = False
        
        if color == 'white':
            paragraph.font.color.rgb = RGBColor(255, 255, 255)
        elif color == 'black':
            paragraph.font.color.rgb = RGBColor(0, 0, 0)
        elif isinstance(color, str) and color.startswith('#'):
            # Handle hex colors like '#CCCCCC'
            hex_color = color.lstrip('#')
            r, g, b = tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
            paragraph.font.color.rgb = RGBColor(r, g, b)
        elif isinstance(color, RGBColor):
            # Handle RGBColor objects
            paragraph.font.color.rgb = color
        else:
            paragraph.font.color.rgb = RGBColor(*self._hex_to_rgb(color or self.text_color))
    
    def apply_caption_style(self, paragraph, color=None, size=12):
        """Small caption text"""
        paragraph.font.name = self.secondary_font
        paragraph.font.size = Pt(size)
        paragraph.font.bold = False
        paragraph.font.color.rgb = RGBColor(*self._hex_to_rgb(color or self.text_color))
    
    def _hex_to_rgb(self, hex_color):
        """Convert hex color to RGB tuple"""
        hex_color = hex_color.lstrip('#')
        if len(hex_color) == 3:
            hex_color = ''.join([c*2 for c in hex_color])
        try:
            return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
        except:
            return (255, 255, 255)  # Default white


class EnhancedPPTXGenerator:
    """
    Enhanced PowerPoint Brand Book Generator with professional layout and styling
    """
    
    def __init__(self):
        self.font_research_agent = FontResearchAgent()
        self.introduction_agent = IntroductionContentAgent()
        self.brand_purpose_agent = BrandPurposeContentAgent()
        self.brand_story_agent = BrandStoryContentAgent()
        self.enhanced_color_agent = EnhancedColorResearchAgent()
        self.iconography_agent = IconographyAgent()
        self.researched_fonts = None
        self.grid = GridLayout()
        self.styler = None
        self.enhanced_color_system = None  # Store the enhanced color system
    
    def _initialize_styling(self, identity_data, company_name=""):
        """Initialize styling system based on brand data"""
        # Get fonts from research
        primary_font = self.researched_fonts['primary_font'] if self.researched_fonts else "Inter"
        secondary_font = self.researched_fonts['secondary_font'] if self.researched_fonts else "Source Sans Pro"
        
        # Get primary color - prefer enhanced color system if available
        if self.enhanced_color_system and self.enhanced_color_system.get('primary_colors'):
            # Use the first primary color from enhanced color system
            first_primary = self.enhanced_color_system['primary_colors'][0]
            primary_color = first_primary['hex']
            print(f"🎨 Using enhanced primary color: {primary_color} ({first_primary['name']})")
        else:
            # Fallback to original palette method
            palette = identity_data.get("palette", {}) if identity_data else {}
            primary_color = self._get_brand_color(palette, "primary", "#2E86AB")
            print(f"🎨 Using fallback primary color: {primary_color}")
        
        text_color = "#000000"  # Black text for white background
        
        self.styler = TextStyler(primary_font, secondary_font, primary_color, text_color)
        
        # Store company name for footers
        self.company_name = company_name
        self.slide_counter = 0
    
    def _update_styling_with_primary_color(self):
        """Update styling to use the first primary color from enhanced color system"""
        if self.enhanced_color_system and self.enhanced_color_system.get('primary_colors') and self.styler:
            first_primary = self.enhanced_color_system['primary_colors'][0]
            new_primary_color = first_primary['hex']
            print(f"🎨 Updating all slide titles to use: {new_primary_color} ({first_primary['name']})")
            
            # Update the styler's primary color
            self.styler.primary_color = new_primary_color
    
    def _get_primary_color_hex(self):
        """Get the primary color hex from enhanced system or fallback"""
        if self.enhanced_color_system and self.enhanced_color_system.get('primary_colors'):
            return self.enhanced_color_system['primary_colors'][0]['hex']
        elif self.styler:
            return self.styler.primary_color
        else:
            return "#2E86AB"  # Safe fallback
    
    def _convert_svg_to_png(self, svg_path: str, size: int = 128) -> str:
        """Convert SVG to PNG for PPTX compatibility"""
        try:
            # Try using cairosvg if available
            try:
                import cairosvg
                png_path = svg_path.replace('.svg', '.png')
                cairosvg.svg2png(
                    url=svg_path, 
                    write_to=png_path,
                    output_width=size,
                    output_height=size
                )
                return png_path
            except ImportError:
                pass
            
            # Try using Pillow with svg2png (wand) if available  
            try:
                from PIL import Image
                from wand.image import Image as WandImage
                
                png_path = svg_path.replace('.svg', '.png')
                with WandImage(filename=svg_path, width=size, height=size) as img:
                    img.format = 'png'
                    img.save(filename=png_path)
                return png_path
            except ImportError:
                pass
            
            # Fallback - return None to trigger placeholder
            print(f"  ⚠️ SVG conversion libraries not available for {svg_path}")
            return None
            
        except Exception as e:
            print(f"  ❌ Error converting SVG to PNG: {e}")
            return None
    
    def _get_industry_icon_categories(self, industry):
        """Get industry-specific icon categories for generation"""
        industry_lower = industry.lower()
        
        # Industry-specific icon mappings
        industry_mappings = {
            "finance": [
                "wallet money icon", "credit card icon", "banking institution icon", 
                "financial chart graph icon", "calculator math icon", "investment growth icon"
            ],
            "fintech": [
                "digital wallet icon", "mobile payment icon", "cryptocurrency coin icon",
                "trading chart icon", "financial analytics icon", "secure payment icon"  
            ],
            "technology": [
                "cloud computing icon", "server database icon", "programming code icon",
                "network connection icon", "artificial intelligence icon", "security lock icon"
            ],
            "healthcare": [
                "medical cross icon", "stethoscope icon", "heart health icon",
                "hospital building icon", "medicine pill icon", "first aid icon"
            ],
            "education": [
                "book learning icon", "graduation cap icon", "school building icon",
                "pencil writing icon", "globe world icon", "certificate diploma icon"
            ],
            "travel": [
                "airplane flight icon", "world map icon", "suitcase luggage icon", 
                "hotel accommodation icon", "camera photography icon", "compass navigation icon"
            ],
            "retail": [
                "shopping bag icon", "shopping cart icon", "store building icon",
                "price tag icon", "delivery package icon", "gift present icon"
            ],
            "food": [
                "restaurant utensils icon", "chef hat icon", "food plate icon",
                "delivery truck icon", "grocery basket icon", "recipe book icon"
            ],
            "real_estate": [
                "house home icon", "building property icon", "key access icon",
                "location pin icon", "contract document icon", "measurement ruler icon"
            ],
            "automotive": [
                "car vehicle icon", "mechanic tools icon", "gas station icon",
                "traffic light icon", "parking icon", "road highway icon"
            ]
        }
        
        # Find matching industry
        for key, categories in industry_mappings.items():
            if key in industry_lower:
                return categories
        
        # Default to technology if no match
        return industry_mappings["technology"]
    
    def _get_brand_color(self, palette, color_name, default):
        """Extract color from palette safely"""
        if not palette:
            return default
        color = palette.get(color_name, default)
        if isinstance(color, list):
            return color[0] if color else default
        return color if color else default
    
    def _get_accessible_text_color(self, palette):
        """Get text color with good contrast against dark backgrounds"""
        primary_color = self._get_brand_color(palette, "primary", "#FFFFFF")
        
        try:
            # Calculate brightness
            hex_color = primary_color.lstrip('#')
            r = int(hex_color[0:2], 16)
            g = int(hex_color[2:4], 16) 
            b = int(hex_color[4:6], 16)
            brightness = (r * 299 + g * 587 + b * 114) / 1000
            
            # Use brand color if bright enough, otherwise white
            return primary_color if brightness > 100 else "#FFFFFF"
        except (ValueError, IndexError):
            return "#FFFFFF"
    
    def _hex_to_rgb(self, hex_color):
        """Convert hex color to RGB tuple"""
        hex_color = hex_color.lstrip('#')
        if len(hex_color) == 3:
            hex_color = ''.join([c*2 for c in hex_color])
        try:
            return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
        except:
            return (255, 255, 255)
    
    def _add_slide_background(self, slide, gradient=False, identity_data=None, bg_color=None):
        """Add background to slide with support for white, black, and gradient backgrounds"""
        if bg_color == 'white':
            # Set solid white background
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(255, 255, 255)  # Pure white
        elif bg_color == 'pitch_black' or bg_color == 'black':
            # Set solid pitch black background
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(0, 0, 0)  # Pure pitch black
        elif gradient and identity_data and identity_data.get("palette"):
            # Create gradient background
            bg_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                0, 0, 
                Inches(10), Inches(7.5)
            )
            palette = identity_data["palette"]
            primary = self._get_brand_color(palette, "primary", "#1a1a1a")
            secondary = self._get_brand_color(palette, "secondary", "#2d2d2d")
            
            fill = bg_shape.fill
            fill.gradient()
            fill.gradient_angle = 135
            gradient = fill.gradient_stops
            gradient[0].color.rgb = RGBColor(*self._hex_to_rgb(primary))
            gradient[1].color.rgb = RGBColor(*self._hex_to_rgb(secondary))
            bg_shape.line.fill.background()
            
            # Move to back
            bg_shape.element.getparent().remove(bg_shape.element)
            slide.shapes._spTree.insert(1, bg_shape.element)
        else:
            # Default: solid white background (changed from dark)
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(255, 255, 255)  # Default to white
    
    def _add_footer(self, slide, slide_number):
        """Add consistent footer with company name and slide number"""
        if not hasattr(self, 'company_name'):
            return
        
        left, top, width, height = self.grid.get_position(0, 7, 12, 1)
        
        # Company name on left
        company_textbox = slide.shapes.add_textbox(left, top, width / 2, height)
        company_frame = company_textbox.text_frame
        company_frame.text = self.company_name
        self.styler.apply_caption_style(company_frame.paragraphs[0], size=10)
        company_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        
        # Slide number on right
        number_textbox = slide.shapes.add_textbox(left + width / 2, top, width / 2, height)
        number_frame = number_textbox.text_frame
        number_frame.text = str(slide_number)
        self.styler.apply_caption_style(number_frame.paragraphs[0], size=10)
        number_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
    
    def _add_logo_to_slide(self, slide, identity_data, col=10, row=0, size=1, opacity=1.0):
        """Add company logo at specified grid position with optional opacity"""
        if not identity_data or not identity_data.get("logos"):
            return
        
        logos = identity_data["logos"]
        logo_path = None
        
        for logo in logos:
            if isinstance(logo, str) and os.path.exists(logo):
                logo_path = logo
                break
        
        if logo_path:
            try:
                left, top, width, height = self.grid.get_position(col, row, size, size)
                logo = slide.shapes.add_picture(logo_path, left, top, width=width, height=height)
                
                # Set logo opacity if specified
                if opacity < 1.0:
                    try:
                        # Access the picture's transparency property
                        logo.element.get_or_add_alpha().val = int(opacity * 100000)  # Alpha in EMU units
                    except Exception as opacity_error:
                        print(f"Could not set logo opacity: {opacity_error}")
                        
            except Exception as e:
                print(f"Could not add logo: {e}")
    
    def _add_geometric_accent(self, slide, identity_data, dark_theme=False):
        """Add subtle geometric accent shape with dark theme support"""
        if not identity_data or not identity_data.get("palette"):
            return
        
        palette = identity_data["palette"]
        
        if dark_theme:
            # Use lighter colors for dark background
            accent_color = self._get_brand_color(palette, "accent", "#FFFFFF")
        else:
            accent_color = self._get_brand_color(palette, "accent", 
                                               self._get_brand_color(palette, "secondary", "#4A90E2"))
        
        # Add accent circle
        left, top, width, height = self.grid.get_position(9, 1, 2, 2)
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
        circle.fill.solid()
        circle.fill.fore_color.rgb = RGBColor(*self._hex_to_rgb(accent_color))
        circle.fill.transparency = 0.7
        circle.line.fill.background()
    
    def _add_logo_as_background(self, slide, identity_data, opacity=0.85):
        """Add logo as full background with specified opacity"""
        if not identity_data or not identity_data.get("logos"):
            return
        
        try:
            # Get the first logo
            logo_path = identity_data["logos"][0]
            
            # Add logo covering the entire slide (10 inches x 7.5 inches standard)
            left = Inches(0)
            top = Inches(0) 
            width = Inches(10)
            height = Inches(7.5)
            
            # Add the logo image
            logo_shape = slide.shapes.add_picture(logo_path, left, top, width=width, height=height)
            
            # Set opacity (transparency)
            try:
                # Method 1: Try to set transparency via fill
                if hasattr(logo_shape, 'fill'):
                    logo_shape.fill.transparency = 1.0 - opacity  # transparency is inverse of opacity
                
                # Method 2: Alternative method for transparency
                elif hasattr(logo_shape, 'element'):
                    # This sets the alpha channel for transparency
                    alpha_value = int(opacity * 100000)  # Convert to percentage for Office
                    
                    # Try to access the picture element and set alpha
                    pic_element = logo_shape.element
                    if pic_element is not None:
                        # This is a more complex approach but might work better
                        pass  # Keep the default opacity for now if complex method needed
                        
            except Exception as e:
                print(f"Could not set logo opacity: {e}")
                # Logo will still be added, just without opacity adjustment
            
            print(f"✅ Added background logo with opacity {opacity}")
            
        except Exception as e:
            print(f"Could not add background logo: {e}")
    
    def _create_title_slide(self, prs, company_name, identity_data, brand_essence):
        """Create professional title slide with black background and enhanced layout"""
        self.slide_counter += 1
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # White background instead of gradient
        self._add_slide_background(slide, gradient=False, identity_data=identity_data, bg_color='white')
        
        # FULL BACKGROUND LOGO with slight opacity reduction
        self._add_logo_as_background(slide, identity_data, opacity=0.85)
        
        # Get dynamic primary color from brand palette
        primary_color = self._get_primary_color_hex()  # Use the enhanced color system
        if identity_data and identity_data.get("palette"):
            palette = identity_data["palette"]
            primary_color = self._get_primary_color_hex()
        
        # RIGHT SIDE: Dynamic colored line above "Brand Identity System"
        line_left, line_top, line_width, line_height = self.grid.get_position(9, 7, 3, 2)
        line_shape = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, 
            line_left, line_top, 
            line_left + line_width, line_top
        )
        # Set line color to brand primary color
        line_rgb = self._hex_to_rgb(primary_color)
        line_shape.line.color.rgb = RGBColor(*line_rgb)
        line_shape.line.width = Pt(0)
        
        # RIGHT SIDE: "Brand Identity System" text centered below line
        left, top, width, height = self.grid.get_position(9, 7, 3, 2)
        brand_system_textbox = slide.shapes.add_textbox(left, top, width, height)
        brand_system_frame = brand_system_textbox.text_frame
        brand_system_frame.text = "Brand Identity System"
        brand_system_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        self.styler.apply_subtitle_style(brand_system_frame.paragraphs[0], size=18, color=primary_color)
        brand_system_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        brand_system_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    def _create_introduction_slide(self, prs, company_name, brand_essence, identity_data, 
                                  industry="", values="", audience=""):
        """Create introduction slide with AI-generated company overview"""
        self.slide_counter += 1
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_slide_background(slide, gradient=False, identity_data=identity_data, bg_color='white')
        
        # Get brand primary color (ensure it uses the agent's chosen primary color)
        primary_color_hex = self._get_primary_color_hex()  # Use enhanced color system
        if identity_data and identity_data.get("palette"):
            palette = identity_data["palette"]
            primary_color_hex = self._get_primary_color_hex()
        
        primary_color_rgb = RGBColor(*self._hex_to_rgb(primary_color_hex))
        
        # Generate AI-powered introduction content
        print("  🤖 Generating AI-powered introduction content...")
        try:
            intro_content = self.introduction_agent.generate_introduction(
                company_name, industry, values, audience, brand_essence
            )
        except Exception as e:
            print(f"  ⚠️ AI introduction generation failed, using fallback: {e}")
            # Fallback to original content
            intro_content = f"The following brand identity system for {company_name}\nis thoughtfully crafted to present our brand in\nan international, engaging, consistent, recognisable\nand proprietary way. Unique in form, versatile in its\napplication and unified by a fundamental principle."
        
        # Content positioned in upper portion - reduced top space and increased width
        left, top, width, height = self.grid.get_position(0.5, 0.8, 10, 4)
        content_textbox = slide.shapes.add_textbox(left, top, width, height)
        content_frame = content_textbox.text_frame
        content_frame.text = intro_content
        content_frame.word_wrap = True
        content_frame.margin_left = 0
        content_frame.margin_right = 0
        content_frame.margin_top = 0
        content_frame.margin_bottom = 0
        
        # Style the content text - black, left-aligned, same as index slide
        for paragraph in content_frame.paragraphs:
            self.styler.apply_body_style(paragraph, color='black', size=20)
            paragraph.alignment = PP_ALIGN.LEFT
            paragraph.space_after = Pt(8)  # Consistent spacing
        
        # Full-width line separator (same width as index slide)
        line_top = Inches(6.5)  # Move line to bottom with title
        line_shape = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(0.5), line_top,
            Inches(9.5), line_top
        )
        line_shape.line.color.rgb = primary_color_rgb
        line_shape.line.width = Pt(2)
        
        # "INTRODUCTION" title in primary color at bottom of slide
        title_top = Inches(6.8)  # Move to bottom
        title_textbox = slide.shapes.add_textbox(
            self.grid.get_position(0.5, 6.5, 1, 0.8)[0], title_top,
            Inches(4), Inches(0.8)
        )
        title_frame = title_textbox.text_frame
        title_frame.text = "INTRODUCTION"
        self.styler.apply_title_style(title_frame.paragraphs[0], size=28, color=primary_color_hex)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        
        
        
       


    def _create_brand_purpose_slide(self, prs, brand_essence, identity_data, 
                                   company_name="", industry="", values="", audience=""):
        """Create brand purpose slide with AI-generated mission and values"""
        self.slide_counter += 1
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_slide_background(slide, gradient=False, identity_data=identity_data, bg_color='white')
        
        # Get brand primary color (ensure it uses the agent's chosen primary color)
        primary_color_hex = self._get_primary_color_hex()  # Use enhanced color system
        if identity_data and identity_data.get("palette"):
            palette = identity_data["palette"]
            primary_color_hex = self._get_primary_color_hex()
        
        primary_color_rgb = RGBColor(*self._hex_to_rgb(primary_color_hex))
        
        # Generate AI-powered brand purpose content
        print("  🤖 Generating AI-powered brand purpose content...")
        try:
            purpose_data = self.brand_purpose_agent.generate_brand_purpose(
                company_name, industry, values, audience, brand_essence
            )
            purpose_content = purpose_data.get("full_content", "")
        except Exception as e:
            print(f"  ⚠️ AI brand purpose generation failed, using fallback: {e}")
            # Fallback content (without "Our Purpose" title)
            purpose_content = ""
            
            if brand_essence and brand_essence.get("brand_positioning"):
                positioning = brand_essence["brand_positioning"]
                if positioning.get("unique_value_proposition"):
                    purpose_content += f"Vision: {positioning['unique_value_proposition']}\n\n"
                if positioning.get("brand_promise"):
                    purpose_content += f"Mission: {positioning['brand_promise']}\n\n"
            
            # Add core values if available
            if brand_essence and brand_essence.get("company_profile", {}).get("core_values"):
                brand_values = brand_essence["company_profile"]["core_values"]
                purpose_content += f"Core Values:\n"
                for value in brand_values:
                    purpose_content += f"• {value}\n"
            else:
                purpose_content += "Core Values:\nOur core values guide everything we do, from innovation to customer service excellence."
        
        # Main content text - positioned in upper area (same as Introduction)
        # Content positioned in upper portion - reduced top space and increased width
        left, top, width, height = self.grid.get_position(0.5, 0.8, 10, 4)
        content_textbox = slide.shapes.add_textbox(left, top, width, height)
        content_frame = content_textbox.text_frame
        content_frame.text = purpose_content
        content_frame.word_wrap = True
        content_frame.margin_left = 0
        content_frame.margin_right = 0
        content_frame.margin_top = 0
        content_frame.margin_bottom = 0
        
        # Style the content text - black, left-aligned, same as introduction slide
        for paragraph in content_frame.paragraphs:
            self.styler.apply_body_style(paragraph, color='black', size=20)
            paragraph.alignment = PP_ALIGN.LEFT
            paragraph.space_after = Pt(8)  # Consistent spacing
        
        # Full-width line separator (same width as introduction slide)
        line_top = Inches(6.5)  # Move line to bottom with title
        line_shape = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(0.5), line_top,
            Inches(9.5), line_top
        )
        line_shape.line.color.rgb = primary_color_rgb
        line_shape.line.width = Pt(2)
        
        # "BRAND PURPOSE" title in primary color at bottom of slide
        title_top = Inches(6.8)  # Move to bottom
        title_textbox = slide.shapes.add_textbox(
            self.grid.get_position(0.5, 6.5, 1, 0.8)[0], title_top,
            Inches(4), Inches(0.8)
        )
        title_frame = title_textbox.text_frame
        title_frame.text = "BRAND PURPOSE"
        self.styler.apply_title_style(title_frame.paragraphs[0], size=28, color=primary_color_hex)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        


    
    def _create_table_of_contents(self, prs, sections, identity_data):
        """Create auto-generated table of contents with INDEX design pattern"""
        self.slide_counter += 1
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # White background
        self._add_slide_background(slide, gradient=False, identity_data=identity_data, bg_color='white')
        
        # Get brand primary color from AI research
        primary_color_hex = self._get_primary_color_hex()  # Use enhanced color system
        if identity_data and identity_data.get("palette"):
            primary_color_hex = self._get_primary_color_hex()
        
        primary_rgb = RGBColor(*self._hex_to_rgb(primary_color_hex))
        
        # Filter sections (remove Table of Contents from list)
        filtered_sections = [s for s in sections if s != "Table of Contents"]
        
        # Map our actual slides to index format with page numbers
        section_data = []
        page_num = 1
        
        for section in filtered_sections:
            section_name = section.replace("1. ", "").replace("2. ", "").replace("3. ", "").replace("4. ", "").replace("5. ", "").replace("6. ", "").replace("7. ", "").replace("8. ", "").replace("9. ", "").replace("10. ", "").replace("11. ", "").replace("12. ", "").replace("13. ", "").replace("14. ", "")
            
            if "Logo" in section_name:
                section_data.append({
                    "name": "Logo Variations",
                    "page": f"{page_num:02d}",
                    "subsections": ["Variation 1", "Variation 2", "Variation 3"]
                })
            elif "Color" in section_name:
                section_data.append({
                    "name": "Brand Colors",
                    "page": f"{page_num:02d}",
                    "subsections": ["Primary Colors", "Secondary Colors", "Color Usage 1/2", "Color Usage 2/2"]
                })
            elif "Typography" in section_name:
                section_data.append({
                    "name": "Typography",
                    "page": f"{page_num:02d}",
                    "subsections": ["Brand Font", "Font Usage", "Text Hierarchy"]
                })
            elif "Imagery & Visuals" in section_name:
                section_data.append({
                    "name": "Imagery & Visuals",
                    "page": f"{page_num:02d}",
                    "subsections": ["Illustrations", "Iconography", "Merchandise"]
                })
            else:
                section_data.append({
                    "name": section_name,
                    "page": f"{page_num:02d}",
                    "subsections": []
                })
            page_num += 1
        
        # Three-column layout
        col1_left = Inches(0.5)
        col2_left = Inches(3.8)
        col3_left = Inches(7.1)
        
        # Split sections into three columns
        items_per_col = len(section_data) // 3 + (1 if len(section_data) % 3 > 0 else 0)
        col1_sections = section_data[:items_per_col]
        col2_sections = section_data[items_per_col:items_per_col*2]
        col3_sections = section_data[items_per_col*2:]
        
        # "INDEX" title at top (brand colored text, bold)
        index_textbox = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5),
            Inches(3), Inches(0.8)
        )
        index_frame = index_textbox.text_frame
        index_frame.text = "INDEX"
        self.styler.apply_title_style(index_frame.paragraphs[0], size=28, color=primary_color_hex)
        index_frame.paragraphs[0].font.bold = True
        
        # Top brand colored line (full width)
        top_line_top = Inches(1.2)
        top_line = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(0.5), top_line_top,
            Inches(9.5), top_line_top
        )
        top_line.line.color.rgb = primary_rgb
        top_line.line.width = Pt(2)
        
        # Create content for each column
        def create_column_content(slide, sections, col_left, start_top=1.7):
            current_top = start_top
            
            for section in sections:
                # Main section with brand colored line and page number
                section_top = Inches(current_top)
                
                # Section name (white text, smaller) - reduced width for closer gap
                section_textbox = slide.shapes.add_textbox(col_left, section_top, 
                                                         Inches(2.0), Inches(0.3))
                section_frame = section_textbox.text_frame
                section_frame.text = section["name"]
                section_frame.margin_left = 0
                section_frame.margin_right = 0
                section_frame.margin_top = 0
                section_frame.margin_bottom = 0
                self.styler.apply_body_style(section_frame.paragraphs[0], size=12, color='black')
                section_frame.paragraphs[0].font.bold = True
                
                # Brand colored accent line (moved closer)
                line_left = col_left + Inches(2.05)
                line_top = section_top + Inches(0.12)
                line_shape = slide.shapes.add_connector(
                    MSO_CONNECTOR.STRAIGHT,
                    line_left, line_top,
                    line_left + Inches(0.4), line_top
                )
                line_shape.line.color.rgb = primary_rgb
                line_shape.line.width = Pt(1.5)
                
                # Page number (brand colored text, smaller) - moved closer
                page_left = col_left + Inches(2.5)
                page_textbox = slide.shapes.add_textbox(page_left, section_top, 
                                                       Inches(0.3), Inches(0.3))
                page_frame = page_textbox.text_frame
                page_frame.text = section["page"]
                page_frame.margin_left = 0
                page_frame.margin_top = 0
                page_frame.margin_bottom = 0
                self.styler.apply_body_style(page_frame.paragraphs[0], size=12, color=primary_color_hex)
                page_frame.paragraphs[0].font.bold = True
                
                current_top += 0.4
                
                # Add subsections (gray text, indented, smaller)
                for subsection in section["subsections"]:
                    sub_top = Inches(current_top)
                    sub_left = col_left + Inches(0.2)  # Indent subsections
                    
                    sub_textbox = slide.shapes.add_textbox(sub_left, sub_top, 
                                                          Inches(2.0), Inches(0.25))
                    sub_frame = sub_textbox.text_frame
                    sub_frame.text = subsection
                    sub_frame.margin_left = 0
                    sub_frame.margin_top = 0
                    sub_frame.margin_bottom = 0
                    self.styler.apply_body_style(sub_frame.paragraphs[0], size=9, color='#666666')
                    
                    current_top += 0.25
                
                current_top += 0.15  # Smaller space between main sections
        
        # Create all three columns
        if col1_sections:
            create_column_content(slide, col1_sections, col1_left)
        if col2_sections:
            create_column_content(slide, col2_sections, col2_left)
        if col3_sections:
            create_column_content(slide, col3_sections, col3_left)
        
        
        

    
    def _create_company_profile_slide(self, prs, brand_essence, identity_data):
        """Create company profile with key-value layout and value chips"""
        self.slide_counter += 1
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_slide_background(slide)
        
        # Get brand primary color
        primary_color_hex = self._get_primary_color_hex()
        if identity_data and identity_data.get("palette"):
            primary_color_hex = self._get_primary_color_hex()
        
        # Title
        left, top, width, height = self.grid.get_position(1, 0, 10, 1)
        title_textbox = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_textbox.text_frame
        title_frame.text = "Company Profile"
        self.styler.apply_title_style(title_frame.paragraphs[0], size=28, color=primary_color_hex)
        
        profile = brand_essence.get("company_profile", {})
        
        # Left side - Key facts
        left, top, width, height = self.grid.get_position(1, 2, 5, 4)
        facts_textbox = slide.shapes.add_textbox(left, top, width, height)
        facts_frame = facts_textbox.text_frame
        
        facts_content = []
        if profile.get('name'):
            facts_content.append(f"Company: {profile['name']}")
        if profile.get('industry'):
            facts_content.append(f"Industry: {profile['industry']}")
        if profile.get('target_audience'):
            facts_content.append(f"Target Audience: {profile['target_audience']}")
        
        facts_frame.text = "\n\n".join(facts_content)
        facts_frame.word_wrap = True
        for paragraph in facts_frame.paragraphs:
            self.styler.apply_body_style(paragraph, size=16)
        
        # Right side - Core values as chips
        if profile.get('core_values'):
            values = profile['core_values']
            palette = identity_data.get("palette", {}) if identity_data else {}
            accent_color = self._get_brand_color(palette, "accent", "#4A90E2")
            secondary_color = self._get_brand_color(palette, "secondary", "#6C757D")
            
            chip_colors = [accent_color, secondary_color]
            
            for i, value in enumerate(values[:6]):  # Limit to 6 values
                row = 2 + (i // 2)
                col = 7 + (i % 2) * 2
                
                left, top, width, height = self.grid.get_position(col, row, 2, 1)
                
                # Create chip background
                chip_bg = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE, 
                    left, top + Inches(0.1), 
                    width - Inches(0.1), height - Inches(0.2)
                )
                chip_bg.fill.solid()
                color = chip_colors[i % len(chip_colors)]
                chip_bg.fill.fore_color.rgb = RGBColor(*self._hex_to_rgb(color))
                chip_bg.line.fill.background()
                
                # Add text on chip
                chip_textbox = slide.shapes.add_textbox(left, top, width, height)
                chip_frame = chip_textbox.text_frame
                chip_frame.text = value
                chip_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                chip_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                self.styler.apply_body_style(chip_frame.paragraphs[0], "#FFFFFF", size=14)
        
       
    
    def _create_market_analysis_slide(self, prs, brand_essence, identity_data):
        """Create market analysis with organized blocks"""
        self.slide_counter += 1
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_slide_background(slide)
        
        # Get brand primary color
        primary_color_hex = self._get_primary_color_hex()
        if identity_data and identity_data.get("palette"):
            primary_color_hex = self._get_primary_color_hex()
        
        # Title
        left, top, width, height = self.grid.get_position(1, 0, 10, 1)
        title_textbox = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_textbox.text_frame
        title_frame.text = "Market Analysis & Insights"
        self.styler.apply_title_style(title_frame.paragraphs[0], size=28, color=primary_color_hex)
        
        analysis = brand_essence.get("market_analysis", {})
        
        # Industry Trends block
        if analysis.get("industry_trends"):
            left, top, width, height = self.grid.get_position(1, 2, 4, 2)
            trends_textbox = slide.shapes.add_textbox(left, top, width, height)
            trends_frame = trends_textbox.text_frame
            trends_content = ["Industry Trends:"] + [f"• {trend}" for trend in analysis["industry_trends"][:5]]
            trends_frame.text = "\n".join(trends_content)
            trends_frame.word_wrap = True
            for i, paragraph in enumerate(trends_frame.paragraphs):
                if i == 0:
                    self.styler.apply_subtitle_style(paragraph, size=18)
                else:
                    self.styler.apply_body_style(paragraph, size=14)
        
        # Competitors block
        if analysis.get("competitor_insights", {}).get("notable_competitors"):
            left, top, width, height = self.grid.get_position(5, 2, 4, 2)
            competitors_textbox = slide.shapes.add_textbox(left, top, width, height)
            competitors_frame = competitors_textbox.text_frame
            competitors = analysis["competitor_insights"]["notable_competitors"][:6]
            competitors_content = ["Key Competitors:"] + [f"• {comp}" for comp in competitors]
            competitors_frame.text = "\n".join(competitors_content)
            competitors_frame.word_wrap = True
            for i, paragraph in enumerate(competitors_frame.paragraphs):
                if i == 0:
                    self.styler.apply_subtitle_style(paragraph, size=18)
                else:
                    self.styler.apply_body_style(paragraph, size=14)
        
        # Design Styles block
        if analysis.get("design_trends", {}).get("design_styles"):
            left, top, width, height = self.grid.get_position(9, 2, 3, 2)
            styles_textbox = slide.shapes.add_textbox(left, top, width, height)
            styles_frame = styles_textbox.text_frame
            styles = analysis["design_trends"]["design_styles"][:4]
            styles_content = ["Design Styles:"] + [f"• {style}" for style in styles]
            styles_frame.text = "\n".join(styles_content)
            styles_frame.word_wrap = True
            for i, paragraph in enumerate(styles_frame.paragraphs):
                if i == 0:
                    self.styler.apply_subtitle_style(paragraph, size=18)
                else:
                    self.styler.apply_body_style(paragraph, size=14)
        
        self._add_footer(slide, self.slide_counter)
    
    def _create_brand_positioning_slide(self, prs, brand_essence, identity_data):
        """Create brand positioning with grid layout"""
        self.slide_counter += 1
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_slide_background(slide)
        
        # Get brand primary color
        primary_color_hex = self._get_primary_color_hex()
        if identity_data and identity_data.get("palette"):
            primary_color_hex = self._get_primary_color_hex()
        
        # Title
        left, top, width, height = self.grid.get_position(1, 0, 10, 1)
        title_textbox = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_textbox.text_frame
        title_frame.text = "Brand Positioning"
        self.styler.apply_title_style(title_frame.paragraphs[0], size=28, color=primary_color_hex)
        
        positioning = brand_essence.get("brand_positioning", {})
        
        # UVP box (top left)
        if positioning.get('unique_value_proposition'):
            left, top, width, height = self.grid.get_position(1, 2, 6, 2)
            uvp_textbox = slide.shapes.add_textbox(left, top, width, height)
            uvp_frame = uvp_textbox.text_frame
            uvp_frame.text = f"Unique Value Proposition:\n{positioning['unique_value_proposition']}"
            uvp_frame.word_wrap = True
            self.styler.apply_subtitle_style(uvp_frame.paragraphs[0], size=16)
            for paragraph in uvp_frame.paragraphs[1:]:
                self.styler.apply_body_style(paragraph, size=14)
        
        # Brand Promise box (top right)
        if positioning.get('brand_promise'):
            left, top, width, height = self.grid.get_position(7, 2, 5, 2)
            promise_textbox = slide.shapes.add_textbox(left, top, width, height)
            promise_frame = promise_textbox.text_frame
            promise_frame.text = f"Brand Promise:\n{positioning['brand_promise']}"
            promise_frame.word_wrap = True
            self.styler.apply_subtitle_style(promise_frame.paragraphs[0], size=16)
            for paragraph in promise_frame.paragraphs[1:]:
                self.styler.apply_body_style(paragraph, size=14)
        
        # Brand Personality chips (middle)
        if positioning.get('brand_personality'):
            personalities = positioning['brand_personality']
            palette = identity_data.get("palette", {}) if identity_data else {}
            accent_color = self._get_brand_color(palette, "accent", "#4A90E2")
            
            for i, personality in enumerate(personalities[:5]):
                col = 1 + i * 2
                left, top, width, height = self.grid.get_position(col, 4, 2, 1)
                
                # Create chip
                chip_bg = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    left, top + Inches(0.1),
                    width - Inches(0.1), height - Inches(0.2)
                )
                chip_bg.fill.solid()
                chip_bg.fill.fore_color.rgb = RGBColor(*self._hex_to_rgb(accent_color))
                chip_bg.fill.transparency = 0.2
                chip_bg.line.color.rgb = RGBColor(*self._hex_to_rgb(accent_color))
                chip_bg.line.width = Pt(1)
                
                # Add text
                chip_textbox = slide.shapes.add_textbox(left, top, width, height)
                chip_frame = chip_textbox.text_frame
                chip_frame.text = personality
                chip_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                chip_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                self.styler.apply_body_style(chip_frame.paragraphs[0], size=12)
        
        # Competitive Advantage (bottom)
        if positioning.get('competitive_advantage'):
            left, top, width, height = self.grid.get_position(1, 5, 11, 2)
            advantage_textbox = slide.shapes.add_textbox(left, top, width, height)
            advantage_frame = advantage_textbox.text_frame
            advantage_frame.text = f"Competitive Advantage:\n{positioning['competitive_advantage']}"
            advantage_frame.word_wrap = True
            self.styler.apply_subtitle_style(advantage_frame.paragraphs[0], size=16)
            for paragraph in advantage_frame.paragraphs[1:]:
                self.styler.apply_body_style(paragraph, size=14)
        
        self._add_footer(slide, self.slide_counter)
    
    def _create_logo_variations_slide(self, prs, logos, identity_data):
        """Create separate slides for each logo variation with usage guidelines"""
        # Create 3 separate slides for each logo variation
        for i, logo_path in enumerate(logos[:3]):  # Ensure we only use first 3 logos
            if os.path.exists(logo_path):
                self._create_single_logo_variation_slide(prs, logo_path, i+1, identity_data)
    
    def _create_single_logo_variation_slide(self, prs, logo_path, variation_number, identity_data):
        """Create a single logo variation slide with left logo and right guidelines"""
        self.slide_counter += 1
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_slide_background(slide, gradient=False, identity_data=identity_data, bg_color='white')
        
        # Get brand primary color (ensure it uses the agent's chosen primary color)
        primary_color_hex = self._get_primary_color_hex()  # Use enhanced color system
        if identity_data and identity_data.get("palette"):
            palette = identity_data["palette"]
            primary_color_hex = self._get_primary_color_hex()
        
        primary_color_rgb = RGBColor(*self._hex_to_rgb(primary_color_hex))
        
        # Left side: Logo
        logo_left, logo_top, logo_width, logo_height = self.grid.get_position(1, 2, 5, 4)
        
        
        # Add logo image
        slide.shapes.add_picture(logo_path, logo_left, logo_top, width=logo_width, height=logo_height)
        
        # Right side: Usage guidelines
        guidelines_left, guidelines_top, guidelines_width, guidelines_height = self.grid.get_position(7, 1, 5, 6)
        
        # Generate AI-powered usage guidelines
        print(f"  🤖 Generating usage guidelines for Variation {variation_number}...")
        try:
            usage_guidelines = self._generate_logo_usage_guidelines(variation_number, identity_data)
        except Exception as e:
            print(f"  ⚠️ AI usage guidelines generation failed, using fallback: {e}")
            usage_guidelines = f"Primary logo for digital and print applications. Maintain clear space and minimum size requirements."
        
        # Usage guidelines text
        guidelines_textbox = slide.shapes.add_textbox(guidelines_left, guidelines_top, guidelines_width, guidelines_height)
        guidelines_frame = guidelines_textbox.text_frame
        guidelines_frame.text = usage_guidelines
        guidelines_frame.word_wrap = True
        guidelines_frame.margin_left = 0
        guidelines_frame.margin_right = 0
        guidelines_frame.margin_top = 0
        guidelines_frame.margin_bottom = 0
        
        # Style the guidelines text - black, left-aligned, same as introduction slide
        for paragraph in guidelines_frame.paragraphs:
            self.styler.apply_body_style(paragraph, color='black', size=16)
            paragraph.alignment = PP_ALIGN.LEFT
            paragraph.space_after = Pt(8)
        
        # Full-width line separator (same width as introduction slide)
        line_top = Inches(6.5)  # Move line to bottom with title
        line_shape = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(0.5), line_top,
            Inches(9.5), line_top
        )
        line_shape.line.color.rgb = primary_color_rgb
        line_shape.line.width = Pt(2)
        
        # Title in primary color at bottom of slide
        title_top = Inches(6.8)  # Move to bottom
        title_textbox = slide.shapes.add_textbox(
            self.grid.get_position(0.5, 7, 1, 0.8)[0], title_top,
            Inches(6), Inches(0.8)
        )
        title_frame = title_textbox.text_frame
        title_frame.text = f"Logo Variation {variation_number}"
        self.styler.apply_title_style(title_frame.paragraphs[0], size=28, color=primary_color_hex)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        
        
    
    def _generate_logo_usage_guidelines(self, variation_number, identity_data):
        """Generate AI-powered logo usage guidelines (15-20 words)"""
        try:
            # Import OpenAI here to avoid dependency issues
            import openai
            import os
            
            # Set up OpenAI client
            openai.api_key = os.getenv('OPENAI_API_KEY')
            
            # Get brand context
            company_name = identity_data.get('company_name', 'Company')
            industry = identity_data.get('industry', 'business')
            
            # Create variation-specific context
            variation_contexts = {
                1: "primary logo for main brand applications",
                2: "secondary logo for compact spaces",
                3: "simplified logo for small sizes"
            }
            
            context = variation_contexts.get(variation_number, "versatile logo for various applications")
            
            # Generate concise explanation for why someone might choose this logo (10-15 words)
            prompt = f"""Generate a concise explanation for why someone might choose Logo Variation {variation_number} for {company_name}.

Context: {context}

Create exactly 10-15 words explaining the benefits or reasons for choosing this specific logo variation.
Focus on practical advantages like versatility, impact, or specific use cases.

Response should be one clear sentence, exactly 10-15 words."""

            response = openai.ChatCompletion.create(
                model="gpt-3.5-turbo",
                messages=[{"role": "user", "content": prompt}],
                max_tokens=40,
                temperature=0.7
            )
            
            guidelines = response.choices[0].message.content.strip()
            
            # Ensure it's within word limit (10-15 words)
            words = guidelines.split()
            if len(words) > 15:
                guidelines = ' '.join(words[:15])
            elif len(words) < 10:
                # Add fallback if too short
                guidelines += f" Perfect for diverse brand applications."
                words = guidelines.split()
                if len(words) > 15:
                    guidelines = ' '.join(words[:15])
            
            return guidelines
            
        except Exception as e:
            # Fallback explanations (10-15 words each)
            fallback_explanations = {
                1: "Choose for maximum impact and brand recognition in primary applications.",
                2: "Perfect for constrained spaces while maintaining professional brand presence and clarity.",
                3: "Ideal for small sizes, watermarks, and applications requiring simplified brand representation."
            }
            return fallback_explanations.get(variation_number, "Versatile option offering flexibility across diverse brand applications and formats.")
    
    def _create_color_palette_slide(self, prs, palette, identity_data, company_name="", industry="", values="", audience=""):
        """Create 4 separate color slides using pre-generated color system"""
        print("🎨 Creating color slides with pre-generated color system...")
        
        # Use the enhanced color system that was generated early in the process
        if self.enhanced_color_system:
            color_system = self.enhanced_color_system
        else:
            # Fallback if somehow not generated earlier
            print("⚠️ Color system not found, generating fallback...")
            color_system = self._generate_fallback_color_system(palette)
        
        # Create 4 separate slides (Primary Colors, Secondary Colors, Primary Usage, Secondary Usage)
        self._create_primary_colors_slide(prs, color_system['primary_colors'], identity_data)
        self._create_secondary_colors_slide(prs, color_system['secondary_colors'], identity_data)  
        self._create_primary_usage_slide(prs, color_system['primary_colors'], identity_data)
        self._create_secondary_usage_slide(prs, color_system['secondary_colors'], identity_data)
    
    def _create_primary_colors_slide(self, prs, primary_colors, identity_data):
        """Create Primary Colors slide with 3-4 colors and detailed specifications"""
        self.slide_counter += 1
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_slide_background(slide, gradient=False, identity_data=identity_data, bg_color='white')
        
        # Get brand primary color for styling
        primary_color_hex = self._get_primary_color_hex()  # Use enhanced color system
        if identity_data and identity_data.get("palette"):
            palette = identity_data["palette"]
            primary_color_hex = self._get_primary_color_hex()
        
        primary_color_rgb = RGBColor(*self._hex_to_rgb(primary_color_hex))
        
        # Color swatches in grid layout - moved upward for better balance
        col = 1
        row = 1.6  # Moved up from 1.5 to 1
        colors_per_row = 2
        
        for i, color_spec in enumerate(primary_colors[:4]):  # Maximum 4 primary colors
            left, top, width, height = self.grid.get_position(
                col + (i % colors_per_row) * 5, 
                row + (i // colors_per_row) * 2.8,  # Reduced from 3.5 to 2.8 for less spacing
                4, 3
            )
            
            # Extract hex color for swatch
            hex_color = color_spec['hex']
            rgb = self._hex_to_rgb(hex_color)
            
            # Color swatch
            swatch = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, left, top, Inches(2), Inches(1.5)
            )
            swatch.fill.solid()
            swatch.fill.fore_color.rgb = RGBColor(*rgb)
            swatch.line.color.rgb = RGBColor(255, 255, 255)
            swatch.line.width = Pt(2)
            
            # Color specifications - reduced gap from color swatch
            spec_left = left + Inches(2.2)  # Reduced from 2.5 to 2.2 for tighter spacing
            spec_textbox = slide.shapes.add_textbox(spec_left, top, Inches(3.3), Inches(1.5))  # Increased width slightly
            spec_frame = spec_textbox.text_frame
            spec_frame.text = f"{color_spec['name']}\n{color_spec['hex']}\n{color_spec['rgb']}\n{color_spec['cmyk']}\n{color_spec['text_recommendation']}"
            spec_frame.word_wrap = True
            spec_frame.margin_left = 0
            spec_frame.margin_right = 0
            spec_frame.margin_top = 0
            spec_frame.margin_bottom = 0
            
            # Style the specifications text
            for paragraph in spec_frame.paragraphs:
                self.styler.apply_body_style(paragraph, color='black', size=10)  # Reduced from 12 to 10
                paragraph.alignment = PP_ALIGN.LEFT
                paragraph.space_after = Pt(3)  # Reduced spacing from 4 to 3
        
        # Full-width line separator
        line_top = Inches(6.5)  # Move line to bottom with title
        line_shape = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(0.5), line_top,
            Inches(9.5), line_top
        )
        line_shape.line.color.rgb = primary_color_rgb
        line_shape.line.width = Pt(2)
        
        # Title in primary color below the line
        title_top = Inches(6.8)  # Move to bottom
        title_textbox = slide.shapes.add_textbox(
            Inches(0.5), title_top,
            Inches(6), Inches(0.8)
        )
        title_frame = title_textbox.text_frame
        title_frame.text = "Primary Colors"
        self.styler.apply_title_style(title_frame.paragraphs[0], size=28, color=primary_color_hex)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        
  
    
    def _create_secondary_colors_slide(self, prs, secondary_colors, identity_data):
        """Create Secondary Colors slide with 5 colors"""
        self.slide_counter += 1
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_slide_background(slide, gradient=False, identity_data=identity_data, bg_color='white')
        
        # Get brand primary color for styling
        primary_color_hex = self._get_primary_color_hex()  # Use enhanced color system
        if identity_data and identity_data.get("palette"):
            palette = identity_data["palette"]
            primary_color_hex = self._get_primary_color_hex()
        
        primary_color_rgb = RGBColor(*self._hex_to_rgb(primary_color_hex))
        
        # Color swatches in grid layout (5 colors, 3 on top, 2 on bottom)
        positions = [
            (1, 1.5), (4, 1.5), (7, 1.5),  # Top row
            (2.5, 4), (5.5, 4)  # Bottom row
        ]
        
        for i, color_spec in enumerate(secondary_colors[:5]):  # Maximum 5 secondary colors
            if i < len(positions):
                col, row = positions[i]
                left, top, width, height = self.grid.get_position(col, row, 3, 2)
                
                # Extract hex color for swatch
                hex_color = color_spec['hex']
                rgb = self._hex_to_rgb(hex_color)
                
                # Color swatch
                swatch = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, left, top, Inches(2), Inches(1)
                )
                swatch.fill.solid()
                swatch.fill.fore_color.rgb = RGBColor(*rgb)
                swatch.line.color.rgb = RGBColor(255, 255, 255)
                swatch.line.width = Pt(2)
                
                # Color specifications below swatch
                spec_top = top + Inches(1.2)
                spec_textbox = slide.shapes.add_textbox(left, spec_top, Inches(2), Inches(1))
                spec_frame = spec_textbox.text_frame
                spec_frame.text = f"{color_spec['name']}\n{color_spec['hex']}\n{color_spec['text_recommendation']}"
                spec_frame.word_wrap = True
                spec_frame.margin_left = 0
                spec_frame.margin_right = 0
                spec_frame.margin_top = 0
                spec_frame.margin_bottom = 0
                
                # Style the specifications text
                for paragraph in spec_frame.paragraphs:
                    self.styler.apply_body_style(paragraph, color='black', size=10)
                    paragraph.alignment = PP_ALIGN.CENTER
                    paragraph.space_after = Pt(3)
        
        # Full-width line separator
        line_top = Inches(6.5)  # Move line to bottom with title
        line_shape = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(0.5), line_top,
            Inches(9.5), line_top
        )
        line_shape.line.color.rgb = primary_color_rgb
        line_shape.line.width = Pt(2)
        
        # Title in primary color below the line
        title_top = Inches(6.8)  # Move to bottom
        title_textbox = slide.shapes.add_textbox(
            Inches(0.5), title_top,
            Inches(6), Inches(0.8)
        )
        title_frame = title_textbox.text_frame
        title_frame.text = "Secondary Colors"
        self.styler.apply_title_style(title_frame.paragraphs[0], size=28, color=primary_color_hex)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        
       
    
    def _create_primary_usage_slide(self, prs, primary_colors, identity_data):
        """Create Primary Color Usage slide (1/2) with usage guidelines for primary colors"""
        self.slide_counter += 1
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_slide_background(slide, gradient=False, identity_data=identity_data, bg_color='white')
        
        # Get brand primary color for styling
        primary_color_hex = self._get_primary_color_hex()
        primary_color_rgb = RGBColor(*self._hex_to_rgb(primary_color_hex))
        
        # Generate primary color usage guidelines
        primary_usage_guidelines = self._generate_primary_color_usage(primary_colors)
        
        # Usage guidelines content in upper portion
        left, top, width, height = self.grid.get_position(0.5, 0.8, 10, 5)
        content_textbox = slide.shapes.add_textbox(left, top, width, height)
        content_frame = content_textbox.text_frame
        content_frame.text = primary_usage_guidelines
        content_frame.word_wrap = True
        content_frame.margin_left = 0
        content_frame.margin_right = 0
        content_frame.margin_top = 0
        content_frame.margin_bottom = 0
        
        # Style the content text - black, left-aligned, smaller size for better fit
        for paragraph in content_frame.paragraphs:
            self.styler.apply_body_style(paragraph, color='black', size=14)  # Reduced from 16 to 14
            paragraph.alignment = PP_ALIGN.LEFT
            paragraph.space_after = Pt(6)  # Reduced from 8 to 6
        
        # Full-width line separator
        line_top = Inches(6.5)  # Move line to bottom with title
        line_shape = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(0.5), line_top,
            Inches(9.5), line_top
        )
        line_shape.line.color.rgb = primary_color_rgb
        line_shape.line.width = Pt(2)
        
        # Title in primary color below the line
        title_top = Inches(6.8)  # Move to bottom
        title_textbox = slide.shapes.add_textbox(
            Inches(0.5), title_top,
            Inches(6), Inches(0.8)
        )
        title_frame = title_textbox.text_frame
        title_frame.text = "Color Usage 1/2"
        self.styler.apply_title_style(title_frame.paragraphs[0], size=28, color=primary_color_hex)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        
       
    
    def _create_secondary_usage_slide(self, prs, secondary_colors, identity_data):
        """Create Secondary Color Usage slide (2/2) with usage guidelines for secondary colors"""
        self.slide_counter += 1
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_slide_background(slide, gradient=False, identity_data=identity_data, bg_color='white')
        
        # Get brand primary color for styling
        primary_color_hex = self._get_primary_color_hex()
        primary_color_rgb = RGBColor(*self._hex_to_rgb(primary_color_hex))
        
        # Generate secondary color usage guidelines
        secondary_usage_guidelines = self._generate_secondary_color_usage(secondary_colors)
        
        # Usage guidelines content in upper portion
        left, top, width, height = self.grid.get_position(0.5, 0.8, 10, 5)
        content_textbox = slide.shapes.add_textbox(left, top, width, height)
        content_frame = content_textbox.text_frame
        content_frame.text = secondary_usage_guidelines
        content_frame.word_wrap = True
        content_frame.margin_left = 0
        content_frame.margin_right = 0
        content_frame.margin_top = 0
        content_frame.margin_bottom = 0
        
        # Style the content text - black, left-aligned, smaller size for better fit
        for paragraph in content_frame.paragraphs:
            self.styler.apply_body_style(paragraph, color='black', size=14)  # Reduced from 16 to 14
            paragraph.alignment = PP_ALIGN.LEFT
            paragraph.space_after = Pt(6)  # Reduced from 8 to 6
        
        # Full-width line separator
        line_top = Inches(6.5)  # Move line to bottom with title
        line_shape = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(0.5), line_top,
            Inches(9.5), line_top
        )
        line_shape.line.color.rgb = primary_color_rgb
        line_shape.line.width = Pt(2)
        
        # Title in primary color below the line
        title_top = Inches(6.8)  # Move to bottom
        title_textbox = slide.shapes.add_textbox(
            Inches(0.5), title_top,
            Inches(6), Inches(0.8)
        )
        title_frame = title_textbox.text_frame
        title_frame.text = "Color Usage 2/2"
        self.styler.apply_title_style(title_frame.paragraphs[0], size=28, color=primary_color_hex)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        
       
    
    def _generate_primary_color_usage(self, primary_colors):
        """Generate shorter usage guidelines for primary colors"""
        color_names = [color['name'] for color in primary_colors]
        
        usage_text = f"""Primary Color Usage

Use for brand-critical elements:
• Logos and main brand marks
• Headlines and key text
• CTA buttons and navigation
• Hero sections

Colors: {', '.join(color_names)}

Best Practices:
- Maximum brand recognition
- Consistent application
- Sufficient contrast
- Accessibility compliance"""
        
        return usage_text
    
    def _generate_secondary_color_usage(self, secondary_colors):
        """Generate shorter usage guidelines for secondary colors"""
        color_names = [color['name'] for color in secondary_colors]
        
        usage_text = f"""Secondary Color Usage

Use for supporting elements:
• Backgrounds and subtle accents
• Secondary text
• Borders and dividers
• Hover states and icons

Colors: {', '.join(color_names)}

Best Practices:
- Create visual hierarchy
- Support, don't compete
- Maintain readability
- Use sparingly"""
        
        return usage_text
    
    def _generate_fallback_color_system(self, original_palette):
        """Generate fallback color system when AI fails"""
        # Convert existing palette to new format
        primary_colors = []
        secondary_colors = []
        
        # Extract first 3 colors as primary
        for i, (color_name, hex_code) in enumerate(list(original_palette.items())[:3]):
            if isinstance(hex_code, list):
                hex_code = hex_code[0] if hex_code else "#CCCCCC"
            
            hex_code = hex_code if hex_code.startswith('#') else f"#{hex_code}"
            rgb = self._hex_to_rgb(hex_code)
            
            primary_colors.append({
                'name': f"Primary {color_name.title()}",
                'hex': hex_code,
                'rgb': f"RGB({rgb[0]}, {rgb[1]}, {rgb[2]})",
                'cmyk': self._rgb_to_cmyk_string(rgb),
                'text_recommendation': self._get_text_recommendation(rgb)
            })
        
        # Add standard secondary colors
        secondary_palette = [
            {'name': 'Light Background', 'hex': '#F8F9FA'},
            {'name': 'Medium Gray', 'hex': '#6B7280'},
            {'name': 'Dark Text', 'hex': '#1F2937'},
            {'name': 'Accent Light', 'hex': '#E5E7EB'},
            {'name': 'Warning Orange', 'hex': '#F59E0B'}
        ]
        
        for color in secondary_palette:
            rgb = self._hex_to_rgb(color['hex'])
            secondary_colors.append({
                'name': color['name'],
                'hex': color['hex'],
                'rgb': f"RGB({rgb[0]}, {rgb[1]}, {rgb[2]})",
                'cmyk': self._rgb_to_cmyk_string(rgb),
                'text_recommendation': self._get_text_recommendation(rgb)
            })
        
        return {
            'primary_colors': primary_colors,
            'secondary_colors': secondary_colors,
            'usage_guidelines': "Use primary colors for main brand elements and headlines. Secondary colors provide supporting elements and backgrounds. Maintain consistent usage across all brand materials."
        }
    
    def _rgb_to_cmyk_string(self, rgb):
        """Convert RGB to CMYK string format"""
        r, g, b = [x/255.0 for x in rgb]
        k = 1 - max(r, g, b)
        if k == 1:
            return "CMYK(0, 0, 0, 100)"
        
        c = (1 - r - k) / (1 - k)
        m = (1 - g - k) / (1 - k)
        y = (1 - b - k) / (1 - k)
        
        return f"CMYK({int(c*100)}, {int(m*100)}, {int(y*100)}, {int(k*100)})"
    
    def _get_text_recommendation(self, rgb):
        """Get text color recommendation based on background color"""
        brightness = sum(rgb) / 3
        return "Use Light Text" if brightness < 128 else "Use Dark Text"
    
    def _create_typography_slide(self, prs, typography, identity_data):
        """Create typography slide with same design as Introduction slide"""
        self.slide_counter += 1
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_slide_background(slide, gradient=False, identity_data=identity_data, bg_color='white')
        
        # Get brand primary color (ensure it uses the agent's chosen primary color)
        primary_color_hex = self._get_primary_color_hex()  # Use enhanced color system
        if identity_data and identity_data.get("palette"):
            palette = identity_data["palette"]
            primary_color_hex = self._get_primary_color_hex()
        
        primary_color_rgb = RGBColor(*self._hex_to_rgb(primary_color_hex))
        
        # Get font information
        primary_font = self.styler.primary_font
        secondary_font = self.styler.secondary_font
        
        # Create typography content matching Introduction slide format
        typography_content = f"Primary Font: {primary_font}\n\nSecondary Font: {secondary_font}\n\nFont Hierarchy:\n• Headlines: {primary_font} - Bold, 28-32pt\n• Subheadlines: {primary_font} - Bold, 20-24pt\n• Body Text: {secondary_font} - Regular, 16-18pt\n• Caption Text: {secondary_font} - Regular, 12-14pt"
        
        # Main content text - positioned in upper area (same as Introduction)
        left, top, width, height = self.grid.get_position(0.5, 0.8, 10, 4)
        content_textbox = slide.shapes.add_textbox(left, top, width, height)
        content_frame = content_textbox.text_frame
        content_frame.text = typography_content
        content_frame.word_wrap = True
        content_frame.margin_left = 0
        content_frame.margin_right = 0
        content_frame.margin_top = 0
        content_frame.margin_bottom = 0
        
        # Style the content text - black, left-aligned, same as introduction slide
        for paragraph in content_frame.paragraphs:
            self.styler.apply_body_style(paragraph, color='black', size=20)
            paragraph.alignment = PP_ALIGN.LEFT
            paragraph.space_after = Pt(8)  # Consistent spacing
        
        # Full-width line separator (same width as introduction slide)
        line_top = Inches(6.5)  # Move line to bottom with title
        line_shape = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(0.5), line_top,
            Inches(9.5), line_top
        )
        line_shape.line.color.rgb = primary_color_rgb
        line_shape.line.width = Pt(2)
        
        # "TYPOGRAPHY" title in primary color below the line (same as Introduction)
        title_top = Inches(6.8)  # Move to bottom
        title_textbox = slide.shapes.add_textbox(
            self.grid.get_position(0.5, 6.5, 1, 0.8)[0], title_top,
            Inches(4), Inches(0.8)
        )
        title_frame = title_textbox.text_frame
        title_frame.text = "TYPOGRAPHY"
        self.styler.apply_title_style(title_frame.paragraphs[0], size=28, color=primary_color_hex)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    
    def _create_icons_display_slide(self, prs, iconography_data, identity_data, company_name="Brand"):
        """Create slide to display generated icons"""
        self.slide_counter += 1
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_slide_background(slide, gradient=False, identity_data=identity_data, bg_color='white')
        
        # Get brand primary color
        primary_color_hex = self._get_primary_color_hex()
        if identity_data and identity_data.get("palette"):
            palette = identity_data["palette"]
            primary_color_hex = self._get_primary_color_hex()
        
        primary_color_rgb = RGBColor(*self._hex_to_rgb(primary_color_hex))
        
        # Simple title text with company name (avoiding duplication)
        if "Icons" in company_name:
            intro_content = company_name
        else:
            intro_content = f"{company_name} Icons"
        
        # Main content text - positioned in upper area
        left, top, width, height = self.grid.get_position(0.5, 0.8, 10, 1)
        content_textbox = slide.shapes.add_textbox(left, top, width, height)
        content_frame = content_textbox.text_frame
        content_frame.text = intro_content
        content_frame.word_wrap = True
        content_frame.margin_left = 0
        content_frame.margin_right = 0
        content_frame.margin_top = 0
        content_frame.margin_bottom = 0
        
        # Style the content text
        for paragraph in content_frame.paragraphs:
            self.styler.apply_body_style(paragraph, color='black', size=20)
            paragraph.alignment = PP_ALIGN.LEFT
            paragraph.space_after = Pt(8)
        
        # Display generated icons if available (15 icons in 5x3 grid)
        if iconography_data and iconography_data.get("icon_generation", {}).get("generated_icons"):
            icons = iconography_data["icon_generation"]["generated_icons"]
            # Handle both old AI format and new Iconify format
            successful_icons = []
            for icon in icons:
                if isinstance(icon, dict):
                    # New Iconify format
                    if icon.get("path"):
                        successful_icons.append(icon)
                    # Old AI format
                    elif "error" not in icon and icon.get("local_path") and icon["local_path"] not in ["Failed to download", "Error"]:
                        successful_icons.append(icon)
            
            if successful_icons:
                # 3 rows x 5 columns layout for 15 icons (1-5, 6-10, 11-15)
                icons_per_row = 5
                max_rows = 3  
                max_icons = min(15, len(successful_icons))
                
                # Use full slide width with minimal margins - compact layout
                slide_width = Inches(10)
                margin = Inches(0.3)  # Very minimal margin
                available_width = slide_width - (2 * margin)
                icon_size = Inches(0.9)  # Smaller icons
                spacing_x = available_width / icons_per_row  # Even distribution
                start_left = margin + (spacing_x - icon_size) / 2  # Center icons in columns
                start_top = Inches(1.8)  # Start higher
                spacing_y = Inches(1.2)  # Reduced vertical spacing
                
                for i, icon in enumerate(successful_icons[:max_icons]):  # Limit to 15 icons
                    try:
                        row = i // icons_per_row
                        col = i % icons_per_row
                        icon_left = start_left + (col * spacing_x)
                        icon_top = start_top + (row * spacing_y)
                        
                        # Get icon path (support both formats)
                        icon_path = icon.get("path") or icon.get("local_path")
                        
                        # Add SVG icon (for Iconify) or image (for AI)
                        if icon_path and os.path.exists(icon_path):
                            if icon_path.endswith('.svg'):
                                # Convert SVG to PNG for PPTX compatibility
                                png_path = self._convert_svg_to_png(icon_path, int(icon_size.inches * 96))
                                if png_path and os.path.exists(png_path):
                                    slide.shapes.add_picture(
                                        png_path,
                                        icon_left, icon_top,
                                        icon_size, icon_size
                                    )
                                else:
                                    # SVG conversion failed - create colored shape placeholder
                                    shape = slide.shapes.add_shape(
                                        MSO_SHAPE.RECTANGLE,
                                        icon_left, icon_top,
                                        icon_size, icon_size
                                    )
                                    # Apply logo color to shape
                                    primary_color_hex = self._get_primary_color_hex()
                                    primary_color_rgb = RGBColor(*self._hex_to_rgb(primary_color_hex))
                                    shape.fill.solid()
                                    shape.fill.fore_color.rgb = primary_color_rgb
                                    shape.line.color.rgb = primary_color_rgb
                            else:
                                # Handle regular image files
                                slide.shapes.add_picture(
                                    icon_path,
                                    icon_left, icon_top,
                                    icon_size, icon_size
                                )
                        
                        # No labels - clean icon-only layout
                        
                    except Exception as e:
                        print(f"  ⚠️ Error adding icon {i+1}: {e}")
                        continue
        
        # Full-width line separator
        line_top = Inches(6.5)  # Move line to bottom with title
        line_shape = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(0.5), line_top,
            Inches(9.5), line_top
        )
        line_shape.line.color.rgb = primary_color_rgb
        line_shape.line.width = Pt(2)
        
        # "ICONOGRAPHY" title
        title_top = Inches(6.8)  # Move to bottom
        title_textbox = slide.shapes.add_textbox(
            self.grid.get_position(0.5, 6.8, 1, 0.8)[0], title_top,
            Inches(4), Inches(0.8)
        )
        title_frame = title_textbox.text_frame
        title_frame.text = "ICONOGRAPHY"
        self.styler.apply_title_style(title_frame.paragraphs[0], size=28, color=primary_color_hex)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

    def _create_iconography_guidelines_slide(self, prs, iconography_data, identity_data):
        """Create slide for iconography usage guidelines"""
        self.slide_counter += 1
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_slide_background(slide, gradient=False, identity_data=identity_data, bg_color='white')
        
        # Get brand primary color
        primary_color_hex = self._get_primary_color_hex()
        if identity_data and identity_data.get("palette"):
            palette = identity_data["palette"]
            primary_color_hex = self._get_primary_color_hex()
        
        primary_color_rgb = RGBColor(*self._hex_to_rgb(primary_color_hex))
        
        # Create concise guidelines content
        if iconography_data and iconography_data.get("system_overview"):
            system_overview = iconography_data["system_overview"]
            style_approach = system_overview.get("style_approach", "minimalist")
            
            guidelines_content = (
                f"Style: {style_approach.title()} design\n"
                f"Color: {primary_color_hex} on black background\n\n"
                f"Size Requirements:\n"
                f"• Minimum: 16px web, 0.5\" print\n"
                f"• Recommended: 24px+ web, 0.75\"+ print\n\n"
                f"Spacing:\n"
                f"• Clear space: 50% of icon width\n"
                f"• Consistent grid alignment"
            )
        else:
            guidelines_content = (
                f"Style: Minimalist design\n"
                f"Color: {primary_color_hex} on black background\n\n"
                f"Size Requirements:\n"
                f"• Minimum: 16px web, 0.5\" print\n"
                f"• Recommended: 24px+ web, 0.75\"+ print\n\n"
                f"Spacing:\n"
                f"• Clear space: 50% of icon width\n"
                f"• Consistent grid alignment"
            )
        
        # Main content text - positioned in upper area
        left, top, width, height = self.grid.get_position(0.5, 0.8, 10, 3)
        content_textbox = slide.shapes.add_textbox(left, top, width, height)
        content_frame = content_textbox.text_frame
        content_frame.text = guidelines_content
        content_frame.word_wrap = True
        content_frame.margin_left = 0
        content_frame.margin_right = 0
        content_frame.margin_top = 0
        content_frame.margin_bottom = 0
        
        # Style the content text
        for paragraph in content_frame.paragraphs:
            self.styler.apply_body_style(paragraph, color='black', size=18)
            paragraph.alignment = PP_ALIGN.LEFT
            paragraph.space_after = Pt(6)
        
        # Full-width line separator
        line_top = Inches(6.5)  # Move line to bottom with title
        line_shape = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(0.5), line_top,
            Inches(9.5), line_top
        )
        line_shape.line.color.rgb = primary_color_rgb
        line_shape.line.width = Pt(2)
        
        # "USAGE GUIDELINES" title
        title_top = Inches(6.8)  # Move to bottom
        title_textbox = slide.shapes.add_textbox(
            self.grid.get_position(0.5, 6.8, 1, 0.8)[0], title_top,
            Inches(6), Inches(0.8)
        )
        title_frame = title_textbox.text_frame
        title_frame.text = "USAGE GUIDELINES"
        self.styler.apply_title_style(title_frame.paragraphs[0], size=28, color=primary_color_hex)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    
    def _create_merchandise_slide_with_logo(self, prs, company_name, logo_file_path, identity_data=None):
        """Create slide displaying brand merchandise mockups using the exact logo from first slide"""
        print("  👕 Generating brand merchandise with exact logo using Flux Pro Kontext...")
        
        try:
            # Get primary color for styling
            primary_color_hex = identity_data.get('primary_color_hex', '#000000') if identity_data else '#000000'
            
            self.slide_counter += 1
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            self._add_slide_background(slide, gradient=False, identity_data=identity_data, bg_color='white')
            
            primary_color_rgb = RGBColor(*self._hex_to_rgb(primary_color_hex))
            
            # Generate t-shirt mockups using Flux Pro with logo reference
            from tools.fal_image_tool import generate_brand_merchandise_with_logo
            
            merchandise_data = generate_brand_merchandise_with_logo(
                company_name=company_name,
                logo_image_path=logo_file_path,
                brand_color=primary_color_hex
            )
            
            if merchandise_data and merchandise_data.get('merchandise_images'):
                merchandise_images = merchandise_data['merchandise_images']
                
                # Layout for 3 t-shirts with gaps and smaller size
                shirt_width = Inches(2.5)  # Smaller width
                shirt_height = Inches(3.2)  # Smaller height  
                gap = Inches(0.8)  # Gap between images
                total_width = (shirt_width * 3) + (gap * 2)  # Total width needed
                start_x = (Inches(10) - total_width) / 2  # Center the images
                start_y = Inches(2.5)  # Top position
                
                shirt_order = ["White", "Black", "Yellow"]  # Consistent order
                
                for i, color_name in enumerate(shirt_order):
                    if color_name in merchandise_images:
                        try:
                            x_pos = start_x + (shirt_width + gap) * i
                            
                            # Add t-shirt image with gaps
                            slide.shapes.add_picture(
                                merchandise_images[color_name], 
                                x_pos, start_y, 
                                width=shirt_width, 
                                height=shirt_height
                            )
                            
                        except Exception as e:
                            print(f"    ⚠️ Failed to add {color_name} t-shirt: {e}")
                            continue
                
                # Full-width line separator (same as Introduction slide)
                line_top = Inches(6.5)  # Move line to bottom with title
                line_shape = slide.shapes.add_connector(
                    MSO_CONNECTOR.STRAIGHT,
                    Inches(0.5), line_top,
                    Inches(9.5), line_top
                )
                line_shape.line.color.rgb = primary_color_rgb
                line_shape.line.width = Pt(2)
                
                # "BRAND APPAREL" title below the line (same as Introduction pattern)
                title_top = Inches(6.8)  # Move to bottom
                title_textbox = slide.shapes.add_textbox(
                    self.grid.get_position(0.5, 6.8, 1, 0.8)[0], title_top,
                    Inches(6), Inches(0.8)
                )
                title_frame = title_textbox.text_frame
                title_frame.text = "BRAND APPAREL"
                self.styler.apply_title_style(title_frame.paragraphs[0], size=28, color=primary_color_hex)
                title_frame.paragraphs[0].font.bold = True
                title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
                
                print(f"  ✅ Brand apparel slide created with {len(merchandise_images)} t-shirt mockups using exact logo")
                
            else:
                # Fallback slide if generation fails
                self._create_merchandise_fallback_slide(prs, company_name, identity_data)
                
        except Exception as e:
            print(f"  ⚠️ Merchandise generation failed: {e}")
            # Create fallback slide
            self._create_merchandise_fallback_slide(prs, company_name, identity_data)
    
    def _create_merchandise_fallback_slide(self, prs, company_name, identity_data=None):
        """Create fallback merchandise slide without generated images"""
        primary_color_hex = identity_data.get('primary_color_hex', '#000000') if identity_data else '#000000'
        
        self.slide_counter += 1
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_slide_background(slide, gradient=False, identity_data=identity_data, bg_color='white')
        
        primary_color_rgb = RGBColor(*self._hex_to_rgb(primary_color_hex))
        
        # Content explaining merchandise concept
        content_textbox = slide.shapes.add_textbox(
            Inches(1), Inches(2),
            Inches(8), Inches(4)
        )
        content_frame = content_textbox.text_frame
        content_frame.text = (
            f"MERCHANDISE APPLICATIONS\n\n"
            f"Brand merchandise serves as powerful marketing tools that extend {company_name}'s "
            f"visual identity into everyday items. Key merchandise categories include:\n\n"
            f"• T-SHIRTS: White, black, and colored variations with consistent logo placement\n"
            f"• APPAREL: Hoodies, caps, and accessories featuring brand colors\n"
            f"• PROMOTIONAL ITEMS: Branded materials for marketing and events\n\n"
            f"All merchandise maintains brand consistency with proper logo usage, "
            f"approved color combinations, and quality standards that reflect {company_name}'s values."
        )
        self.styler.apply_body_style(content_frame.paragraphs[0], size=14, color='#333333')
        content_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        
        # Add placeholder boxes representing t-shirts with gaps and smaller size
        box_width = Inches(2.5)  # Smaller width
        box_height = Inches(1.2)  # Smaller height
        gap = Inches(0.8)  # Gap between boxes
        total_width = (box_width * 3) + (gap * 2)
        start_x = (Inches(10) - total_width) / 2  # Center the boxes
        y_pos = Inches(6.5)
        
        colors = ["#FFFFFF", "#000000", "#FFD700"]  # White, Black, Yellow
        
        for i, bg_color in enumerate(colors):
            x_pos = start_x + (box_width + gap) * i
            
            # Create rectangle representing t-shirt (no labels)
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, x_pos, y_pos, box_width, box_height
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(*self._hex_to_rgb(bg_color))
            shape.line.color.rgb = RGBColor(200, 200, 200)
        
        # Full-width line separator (same as Introduction slide)
        line_top = Inches(6.5)  # Move line to bottom with title
        line_shape = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(0.5), line_top,
            Inches(9.5), line_top
        )
        line_shape.line.color.rgb = primary_color_rgb
        line_shape.line.width = Pt(2)
        
        # "BRAND APPAREL" title below the line (same as Introduction pattern)
        title_top = Inches(6.8)  # Move to bottom
        title_textbox = slide.shapes.add_textbox(
            self.grid.get_position(0.5, 6.8, 1, 0.8)[0], title_top,
            Inches(6), Inches(0.8)
        )
        title_frame = title_textbox.text_frame
        title_frame.text = "BRAND APPAREL"
        self.styler.apply_title_style(title_frame.paragraphs[0], size=28, color=primary_color_hex)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        
        print(f"  ✅ Fallback brand apparel slide created for {company_name}")
    
    def _create_brand_mugs_slide(self, prs, company_name, logo_file_path, identity_data=None):
        """Create slide displaying brand mug mockups using the exact logo from first slide"""
        print("  ☕ Generating brand mugs with exact logo using Flux Pro Kontext...")
        
        try:
            # Get primary color for styling and smart mug color
            primary_color_hex = identity_data.get('primary_color_hex', '#000000') if identity_data else '#000000'
            
            self.slide_counter += 1
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            self._add_slide_background(slide, gradient=False, identity_data=identity_data, bg_color='white')
            
            primary_color_rgb = RGBColor(*self._hex_to_rgb(primary_color_hex))
            
            # Generate mug mockups using Flux Pro with logo reference
            from tools.fal_image_tool import generate_brand_mugs_with_logo
            
            mug_data = generate_brand_mugs_with_logo(
                company_name=company_name,
                logo_image_path=logo_file_path,
                brand_color=primary_color_hex
            )
            
            if mug_data and mug_data.get('mug_images'):
                mug_images = mug_data['mug_images']
                
                # Layout for 3 mugs with gaps and smaller size (same as t-shirts)
                mug_width = Inches(2.5)  # Smaller width
                mug_height = Inches(3.2)  # Smaller height  
                gap = Inches(0.8)  # Gap between images
                total_width = (mug_width * 3) + (gap * 2)  # Total width needed
                start_x = (Inches(10) - total_width) / 2  # Center the images
                start_y = Inches(2.5)  # Top position
                
                mug_order = ["White", "Black", "Primary"]  # Consistent order
                
                for i, color_name in enumerate(mug_order):
                    if color_name in mug_images:
                        try:
                            x_pos = start_x + (mug_width + gap) * i
                            
                            # Add mug image with gaps
                            slide.shapes.add_picture(
                                mug_images[color_name], 
                                x_pos, start_y, 
                                width=mug_width, 
                                height=mug_height
                            )
                            
                        except Exception as e:
                            print(f"    ⚠️ Failed to add {color_name} mug: {e}")
                            continue
                
                # Full-width line separator (same as Introduction slide)
                line_top = Inches(6.5)  # Move line to bottom with title
                line_shape = slide.shapes.add_connector(
                    MSO_CONNECTOR.STRAIGHT,
                    Inches(0.5), line_top,
                    Inches(9.5), line_top
                )
                line_shape.line.color.rgb = primary_color_rgb
                line_shape.line.width = Pt(2)
                
                # "BRAND MUGS" title below the line (same as Introduction pattern)
                title_top = Inches(6.8)  # Move to bottom
                title_textbox = slide.shapes.add_textbox(
                    self.grid.get_position(0.5, 6.8, 1, 0.8)[0], title_top,
                    Inches(6), Inches(0.8)
                )
                title_frame = title_textbox.text_frame
                title_frame.text = "BRAND MUGS"
                self.styler.apply_title_style(title_frame.paragraphs[0], size=28, color=primary_color_hex)
                title_frame.paragraphs[0].font.bold = True
                title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
                
                print(f"  ✅ Brand mugs slide created with {len(mug_images)} mug mockups using exact logo")
                
            else:
                # Fallback slide if generation fails
                self._create_brand_mugs_fallback_slide(prs, company_name, identity_data)
                
        except Exception as e:
            print(f"  ⚠️ Mug generation failed: {e}")
            # Create fallback slide
            self._create_brand_mugs_fallback_slide(prs, company_name, identity_data)
    
    def _create_brand_mugs_fallback_slide(self, prs, company_name, identity_data=None):
        """Create fallback mug slide without generated images"""
        primary_color_hex = identity_data.get('primary_color_hex', '#000000') if identity_data else '#000000'
        
        self.slide_counter += 1
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_slide_background(slide, gradient=False, identity_data=identity_data, bg_color='white')
        
        primary_color_rgb = RGBColor(*self._hex_to_rgb(primary_color_hex))
        
        # Content explaining mug concept
        content_textbox = slide.shapes.add_textbox(
            Inches(1), Inches(2),
            Inches(8), Inches(4)
        )
        content_frame = content_textbox.text_frame
        content_frame.text = (
            f"BRAND MUGS\n\n"
            f"Coffee mugs and drinkware serve as daily brand touchpoints that extend {company_name}'s "
            f"visual identity into workplace and personal environments. Key mug applications include:\n\n"
            f"• OFFICE MUGS: White, black, and brand color variations for corporate settings\n"
            f"• PROMOTIONAL DRINKWARE: Branded mugs for events and client gifts\n"
            f"• RETAIL MERCHANDISE: Consumer products featuring consistent brand elements\n\n"
            f"All drinkware maintains brand consistency with proper logo placement, "
            f"appropriate color combinations, and quality materials that reflect {company_name}'s standards."
        )
        self.styler.apply_body_style(content_frame.paragraphs[0], size=14, color='#333333')
        content_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        
        # Add placeholder boxes representing mugs with gaps and smaller size
        box_width = Inches(2.5)  # Smaller width
        box_height = Inches(1.2)  # Smaller height
        gap = Inches(0.8)  # Gap between boxes
        total_width = (box_width * 3) + (gap * 2)
        start_x = (Inches(10) - total_width) / 2  # Center the boxes
        y_pos = Inches(6.5)
        
        colors = ["#FFFFFF", "#000000", primary_color_hex]  # White, Black, Primary
        
        for i, bg_color in enumerate(colors):
            x_pos = start_x + (box_width + gap) * i
            
            # Create rectangle representing mug (no labels)
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, x_pos, y_pos, box_width, box_height
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(*self._hex_to_rgb(bg_color))
            shape.line.color.rgb = RGBColor(200, 200, 200)
        
        # Full-width line separator (same as Introduction slide)
        line_top = Inches(6.5)  # Move line to bottom with title
        line_shape = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(0.5), line_top,
            Inches(9.5), line_top
        )
        line_shape.line.color.rgb = primary_color_rgb
        line_shape.line.width = Pt(2)
        
        # "BRAND MUGS" title below the line (same as Introduction pattern)
        title_top = Inches(6.8)  # Move to bottom
        title_textbox = slide.shapes.add_textbox(
            self.grid.get_position(0.5, 6.8, 1, 0.8)[0], title_top,
            Inches(6), Inches(0.8)
        )
        title_frame = title_textbox.text_frame
        title_frame.text = "BRAND MUGS"
        self.styler.apply_title_style(title_frame.paragraphs[0], size=28, color=primary_color_hex)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        
        print(f"  ✅ Fallback brand mugs slide created for {company_name}")
    
    
    def _create_brand_illustrations_slide(self, prs, company_name, industry, values, audience, brand_essence="", identity_data=None):
        """Create slide displaying AI-generated brand illustrations with improved layout and organization"""
        print("  🎨 Creating brand illustrations slide...")
        
        try:
            # Always generate fresh professional illustrations using Recraft V3
            print("  🎨 Generating fresh brand illustrations with Recraft V3...")
            from tools.fal_image_tool import generate_professional_brand_illustrations
            
            primary_color_hex = self._get_primary_color_hex()
            
            illustrations_data = generate_professional_brand_illustrations(
                company_name=company_name,
                industry=industry,
                values=values,
                audience=audience,
                brand_essence=brand_essence,
                primary_color=primary_color_hex,
                num_illustrations=4  # Restored to 4 after fixing illustration display issue
            )
            
            successful_illustrations = [
                ill for ill in illustrations_data.get('illustrations', []) 
                if 'error' not in ill and ill.get('local_path') and ill['local_path'] != "Failed to download"
            ]
            
            if successful_illustrations:
                self.slide_counter += 1
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                self._add_slide_background(slide, gradient=False, identity_data=identity_data, bg_color='white')
                
                # Get brand primary color
                primary_color_hex = self._get_primary_color_hex()
                if identity_data and identity_data.get("palette"):
                    primary_color_hex = self._get_primary_color_hex()
                
                primary_color_rgb = RGBColor(*self._hex_to_rgb(primary_color_hex))
                
                # Enhanced grid layout with consistent spacing and framing
                grid_start_top = Inches(0.4)
                illustration_width = Inches(2.6)
                illustration_height = Inches(1.5)
                horizontal_spacing = Inches(0.5)  # Equal spacing between columns
                vertical_spacing = Inches(0.3)    # Consistent vertical spacing
                
                # Use illustrations directly without complex categorization
                for idx, illustration in enumerate(successful_illustrations[:6]):
                    if os.path.exists(illustration['local_path']):
                        # Calculate position in organized 2x3 grid
                        col = idx % 3
                        row = idx // 3
                        
                        # Centered layout with equal padding
                        left = Inches(1.0) + col * (illustration_width + horizontal_spacing)
                        top = grid_start_top + row * (illustration_height + vertical_spacing + Inches(0.4))  # Space for labels
                        
                        try:
                            # DETAILED DEBUGGING - Add illustration image FIRST (so it's not hidden behind frame)
                            print(f"    📷 DEBUGGING: Loading image: {illustration['local_path']}")
                            print(f"    📂 File exists check: {os.path.exists(illustration['local_path'])}")
                            print(f"    📐 Position - Left: {left}, Top: {top}")
                            print(f"    📏 Size - Width: {illustration_width}, Height: {illustration_height}")
                            
                            # Image should be clean PNG from fal_image_tool
                            try:
                                from PIL import Image
                                # Quick validation of clean image
                                with Image.open(illustration['local_path']) as img:
                                    print(f"    ✅ Final validation: {img.format}, {img.size}, {img.mode}")
                            
                                # Add the cleaned image to PowerPoint
                                picture = slide.shapes.add_picture(
                                    illustration['local_path'],
                                    left, top,
                                    width=illustration_width, 
                                    height=illustration_height
                                )
                                
                            except Exception as img_error:
                                print(f"    ❌ Image insertion failed: {img_error}")
                                # Skip this illustration rather than using fallback
                                continue
                            print(f"    ✅ Image loaded successfully!")
                            print(f"    🏷️ Picture name: {picture.name}")
                            print(f"    📊 Picture dimensions: {picture.width} x {picture.height}")
                            print(f"    📍 Picture position: ({picture.left}, {picture.top})")
                            
                            # Add category label above each illustration
                            category_label = f"Visual {idx + 1}"
                            category_top = top - Inches(0.25)
                            category_textbox = slide.shapes.add_textbox(
                                left, category_top,
                                illustration_width, Inches(0.2)
                            )
                            category_frame = category_textbox.text_frame
                            category_frame.text = category_label
                            category_frame.margin_left = 0
                            category_frame.margin_right = 0
                            category_frame.margin_top = 0
                            category_frame.margin_bottom = 0
                            self.styler.apply_body_style(category_frame.paragraphs[0], size=8, color=primary_color_hex)
                            category_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                            category_frame.paragraphs[0].font.bold = True
                            
                            # Add concept label below each illustration
                            label_top = top + illustration_height + Inches(0.08)
                            label_textbox = slide.shapes.add_textbox(
                                left, label_top,
                                illustration_width, Inches(0.25)
                            )
                            label_frame = label_textbox.text_frame
                            concept_short = illustration['concept'][:35] + "..." if len(illustration['concept']) > 35 else illustration['concept']
                            label_frame.text = concept_short
                            label_frame.margin_left = 0
                            label_frame.margin_right = 0
                            label_frame.margin_top = 0
                            label_frame.margin_bottom = 0
                            self.styler.apply_body_style(label_frame.paragraphs[0], size=9, color='#666666')
                            label_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                            
                        except Exception as e:
                            print(f"  ⚠️ Error adding illustration {idx+1}: {e}")
                            continue
                
                # Usage guidelines section
                guidelines_top = grid_start_top + 2 * (illustration_height + vertical_spacing + Inches(0.4)) + Inches(0.3)
                guidelines_textbox = slide.shapes.add_textbox(
                    Inches(1.0), guidelines_top,
                    Inches(8.0), Inches(0.8)
                )
                guidelines_frame = guidelines_textbox.text_frame
                guidelines_frame.text = (
                    "USAGE NOTES: These illustrations define the visual storytelling style of the brand. "
                    "Use them consistently across presentations, reports, and marketing materials to maintain brand coherence."
                )
                guidelines_frame.margin_left = 0
                guidelines_frame.margin_right = 0
                guidelines_frame.margin_top = 0
                guidelines_frame.margin_bottom = 0
                guidelines_frame.word_wrap = True
                self.styler.apply_body_style(guidelines_frame.paragraphs[0], size=10, color='black')
                guidelines_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                
                # Full-width line separator at bottom (matching introduction slide pattern)
                line_top = Inches(6.5)  # Move line to bottom with title
                line_shape = slide.shapes.add_connector(
                    MSO_CONNECTOR.STRAIGHT,
                    Inches(0.5), line_top,
                    Inches(9.5), line_top
                )
                line_shape.line.color.rgb = primary_color_rgb
                line_shape.line.width = Pt(2)
                
                # "BRAND ILLUSTRATIONS" title in primary color below the line
                title_top = Inches(6.8)  # Move to bottom
                title_textbox = slide.shapes.add_textbox(
                    self.grid.get_position(0.5, 6.5, 1, 0.8)[0], title_top,
                    Inches(6), Inches(0.8)
                )
                title_frame = title_textbox.text_frame
                title_frame.text = "BRAND ILLUSTRATIONS"
                title_frame.margin_left = 0
                title_frame.margin_right = 0
                title_frame.margin_top = 0
                title_frame.margin_bottom = 0
                self.styler.apply_title_style(title_frame.paragraphs[0], size=28, color=primary_color_hex)
                title_frame.paragraphs[0].font.bold = True
                title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
                
                print(f"  ✅ Enhanced brand illustrations slide created with {len(successful_illustrations[:6])} illustrations")
                
            else:
                # Fallback slide with no illustrations
                self._create_brand_illustrations_fallback_slide(prs, company_name, identity_data)
                
        except Exception as e:
            print(f"  ⚠️ Brand illustrations generation failed: {e}")
            # Create fallback slide
            self._create_brand_illustrations_fallback_slide(prs, company_name, identity_data)
    
    def _create_brand_illustrations_fallback_slide(self, prs, company_name, identity_data=None):
        """Create fallback illustrations slide without generated images"""
        self.slide_counter += 1
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_slide_background(slide, gradient=False, identity_data=identity_data, bg_color='white')
        
        # Get brand primary color
        primary_color_hex = self._get_primary_color_hex()
        primary_color_rgb = RGBColor(*self._hex_to_rgb(primary_color_hex))
        
        # Content description positioned in upper area (matching introduction slide)
        left, top, width, height = self.grid.get_position(0.5, 0.8, 10, 4)
        content_textbox = slide.shapes.add_textbox(left, top, width, height)
        content_frame = content_textbox.text_frame
        content_frame.text = (
            f"Professional brand illustrations for {company_name} would include:\n\n"
            "• Industry-specific visual concepts\n"
            "• Brand essence representations\n"
            "• Modern design aesthetic\n"
            "• Consistent visual language\n"
            "• Scalable vector graphics\n\n"
            "Style: DESIGN approach with ULTRAMARINE color palette\n"
            "Format: Professional business illustrations suitable for all brand touchpoints"
        )
        content_frame.margin_left = 0
        content_frame.margin_right = 0
        content_frame.margin_top = 0
        content_frame.margin_bottom = 0
        content_frame.word_wrap = True
        
        for paragraph in content_frame.paragraphs:
            self.styler.apply_body_style(paragraph, size=20, color='black')
            paragraph.alignment = PP_ALIGN.LEFT
            paragraph.space_after = Pt(8)
        
        # Full-width line separator (matching introduction slide pattern)
        line_top = Inches(6.5)  # Move line to bottom with title
        line_shape = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(0.5), line_top,
            Inches(9.5), line_top
        )
        line_shape.line.color.rgb = primary_color_rgb
        line_shape.line.width = Pt(2)
        
        # "BRAND ILLUSTRATIONS" title in primary color below the line
        title_top = Inches(6.8)  # Move to bottom
        title_textbox = slide.shapes.add_textbox(
            self.grid.get_position(0.5, 6.5, 1, 0.8)[0], title_top,
            Inches(6), Inches(0.8)
        )
        title_frame = title_textbox.text_frame
        title_frame.text = "BRAND ILLUSTRATIONS"
        title_frame.margin_left = 0
        title_frame.margin_right = 0
        title_frame.margin_top = 0
        title_frame.margin_bottom = 0
        self.styler.apply_title_style(title_frame.paragraphs[0], size=28, color=primary_color_hex)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        
        print("  ✅ Brand illustrations fallback slide created")
    
    def _create_text_slide(self, prs, title, content, max_words=80, identity_data=None):
        """Create text-based slide with pagination if needed"""
        chunks = self._paginate_text(content, max_words)
        
        # Get brand primary color
        primary_color_hex = self._get_primary_color_hex()
        if identity_data and identity_data.get("palette"):
            primary_color_hex = self._get_primary_color_hex()
        
        for idx, chunk in enumerate(chunks):
            self.slide_counter += 1
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            self._add_slide_background(slide)
            
            # Title with pagination indicator if multiple slides
            left, top, width, height = self.grid.get_position(1, 0, 10, 1)
            title_textbox = slide.shapes.add_textbox(left, top, width, height)
            title_frame = title_textbox.text_frame
            slide_title = f"{title} ({idx+1}/{len(chunks)})" if len(chunks) > 1 else title
            title_frame.text = slide_title
            self.styler.apply_title_style(title_frame.paragraphs[0], size=28, color=primary_color_hex)
            
            # Content
            left, top, width, height = self.grid.get_position(1, 2, 10, 4)
            content_textbox = slide.shapes.add_textbox(left, top, width, height)
            content_frame = content_textbox.text_frame
            content_frame.text = chunk
            content_frame.word_wrap = True
            
            for paragraph in content_frame.paragraphs:
                self.styler.apply_body_style(paragraph, size=16)
            
            self._add_footer(slide, self.slide_counter)
    
    def _paginate_text(self, text, words_per_slide=80):
        """Split text into chunks for pagination"""
        if not text or not text.strip():
            return [""]
        
        words = text.split()
        if len(words) <= words_per_slide:
            return [text]
        
        chunks = []
        current_words = []
        
        for word in words:
            current_words.append(word)
            
            if len(current_words) >= words_per_slide:
                if word.endswith('.') or word.endswith('!') or word.endswith('?'):
                    chunks.append(' '.join(current_words))
                    current_words = []
                elif len(current_words) >= words_per_slide + 10:
                    chunks.append(' '.join(current_words))
                    current_words = []
        
        if current_words:
            chunks.append(' '.join(current_words))
        
        return chunks
    
    def _create_brand_purpose_slides(self, prs, brand_essence, identity_data, 
                                    company_name="", industry="", values="", audience=""):
        """Create three separate brand purpose slides: Vision, Mission, Core Values"""
        
        # Generate AI-powered brand purpose content once
        print("  🤖 Generating AI-powered brand purpose content...")
        try:
            purpose_data = self.brand_purpose_agent.generate_brand_purpose(
                company_name, industry, values, audience, brand_essence
            )
            
            # First try to get structured data
            vision_content = purpose_data.get("vision", "")
            mission_content = purpose_data.get("mission", "")
            values_content = purpose_data.get("values_content", "")
            
            # If structured data is missing, parse from full_content
            if not vision_content or not mission_content or not values_content:
                full_content = purpose_data.get("full_content", "")
                if full_content:
                    parsed_vision, parsed_mission, parsed_values = self._parse_brand_purpose_content(full_content)
                    vision_content = parsed_vision or vision_content
                    mission_content = parsed_mission or mission_content  
                    values_content = parsed_values or values_content
            
        except Exception as e:
            print(f"  ⚠️ AI brand purpose generation failed, using fallback: {e}")
            # Fallback content
            vision_content = "To be the preferred partner for businesses seeking excellence and innovation in their field."
            mission_content = "We deliver exceptional solutions that drive growth and create lasting value for our clients and communities."
            values_content = "Our principles guide everything we do, from innovation to customer service excellence."
        
        # Create Vision slide
        self._create_single_purpose_slide(prs, identity_data, "BRAND PURPOSE (VISION)", vision_content)
        
        # Create Mission slide  
        self._create_single_purpose_slide(prs, identity_data, "BRAND PURPOSE (MISSION)", mission_content)
        
        # Create Core Values slide
        self._create_single_purpose_slide(prs, identity_data, "BRAND PURPOSE (VALUES)", values_content)
    
    def _parse_brand_purpose_content(self, full_content):
        """Parse full brand purpose content into Vision, Mission, Core Values"""
        if not full_content:
            return "", "", ""
            
        vision = ""
        mission = ""
        values = ""
        
        lines = full_content.split('\n')
        current_section = None
        current_content = []
        
        for line in lines:
            line = line.strip()
            
            if line.startswith('Vision:'):
                if current_section and current_content:
                    # Save previous section
                    content = ' '.join(current_content).strip()
                    if current_section == "vision":
                        vision = content
                    elif current_section == "mission":
                        mission = content
                    elif current_section == "values":
                        values = content
                
                current_section = "vision"
                current_content = [line.replace('Vision:', '').strip()]
                
            elif line.startswith('Mission:'):
                if current_section and current_content:
                    # Save previous section
                    content = ' '.join(current_content).strip()
                    if current_section == "vision":
                        vision = content
                    elif current_section == "mission":
                        mission = content
                    elif current_section == "values":
                        values = content
                
                current_section = "mission"
                current_content = [line.replace('Mission:', '').strip()]
                
            elif line.startswith('Core Values:'):
                if current_section and current_content:
                    # Save previous section
                    content = ' '.join(current_content).strip()
                    if current_section == "vision":
                        vision = content
                    elif current_section == "mission":
                        mission = content
                    elif current_section == "values":
                        values = content
                
                current_section = "values"
                current_content = [line.replace('Core Values:', '').strip()]
                
            elif line and current_section:
                current_content.append(line)
        
        # Save the last section
        if current_section and current_content:
            content = ' '.join(current_content).strip()
            if current_section == "vision":
                vision = content
            elif current_section == "mission":
                mission = content
            elif current_section == "values":
                values = content
        
        return vision, mission, values
    
    def _create_single_purpose_slide(self, prs, identity_data, title, content):
        """Create a single brand purpose slide with same design as Introduction slide"""
        self.slide_counter += 1
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_slide_background(slide, gradient=False, identity_data=identity_data, bg_color='white')
        
        # Get brand primary color (ensure it uses the agent's chosen primary color)
        primary_color_hex = self._get_primary_color_hex()  # Use enhanced color system
        if identity_data and identity_data.get("palette"):
            palette = identity_data["palette"]
            primary_color_hex = self._get_primary_color_hex()
        
        primary_color_rgb = RGBColor(*self._hex_to_rgb(primary_color_hex))
        
        # Main content text - positioned in upper area (same as Introduction)
        # Content positioned in upper portion - reduced top space and increased width
        left, top, width, height = self.grid.get_position(0.5, 0.8, 10, 4)
        content_textbox = slide.shapes.add_textbox(left, top, width, height)
        content_frame = content_textbox.text_frame
        content_frame.text = content
        content_frame.word_wrap = True
        content_frame.margin_left = 0
        content_frame.margin_right = 0
        content_frame.margin_top = 0
        content_frame.margin_bottom = 0
        
        # Style the content text - black, left-aligned, same as introduction slide
        for paragraph in content_frame.paragraphs:
            self.styler.apply_body_style(paragraph, color='black', size=20)
            paragraph.alignment = PP_ALIGN.LEFT
            paragraph.space_after = Pt(8)  # Consistent spacing
        
        # Full-width line separator (same width as introduction slide)
        line_top = Inches(6.5)  # Move line to bottom with title
        line_shape = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(0.5), line_top,
            Inches(9.5), line_top
        )
        line_shape.line.color.rgb = primary_color_rgb
        line_shape.line.width = Pt(2)
        
        # Title in primary color below the line (same as Introduction)
        title_top = Inches(6.8)  # Move to bottom
        title_textbox = slide.shapes.add_textbox(
            self.grid.get_position(0.5, 6.5, 1, 0.8)[0], title_top,
            Inches(6), Inches(0.8)
        )
        title_frame = title_textbox.text_frame
        title_frame.text = title
        self.styler.apply_title_style(title_frame.paragraphs[0], size=28, color=primary_color_hex)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    
    def _create_brand_story_slide(self, prs, identity_data, title, content):
        """Create brand story slide with same design as Introduction slide"""
        self.slide_counter += 1
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_slide_background(slide, gradient=False, identity_data=identity_data, bg_color='white')
        
        # Get brand primary color (ensure it uses the agent's chosen primary color)
        primary_color_hex = self._get_primary_color_hex()  # Use enhanced color system
        if identity_data and identity_data.get("palette"):
            palette = identity_data["palette"]
            primary_color_hex = self._get_primary_color_hex()
        
        primary_color_rgb = RGBColor(*self._hex_to_rgb(primary_color_hex))
        
        # Main content text - positioned in upper area (same as Introduction)
        # Content positioned in upper portion - reduced top space and increased width
        left, top, width, height = self.grid.get_position(0.5, 0.8, 10, 4)
        content_textbox = slide.shapes.add_textbox(left, top, width, height)
        content_frame = content_textbox.text_frame
        content_frame.text = content
        content_frame.word_wrap = True
        content_frame.margin_left = 0
        content_frame.margin_right = 0
        content_frame.margin_top = 0
        content_frame.margin_bottom = 0
        
        # Style the content text - black, left-aligned, same as introduction slide
        for paragraph in content_frame.paragraphs:
            self.styler.apply_body_style(paragraph, color='black', size=20)
            paragraph.alignment = PP_ALIGN.LEFT
            paragraph.space_after = Pt(8)  # Consistent spacing
        
        # Full-width line separator (same width as introduction slide)
        line_top = Inches(6.5)  # Move line to bottom with title
        line_shape = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(0.5), line_top,
            Inches(9.5), line_top
        )
        line_shape.line.color.rgb = primary_color_rgb
        line_shape.line.width = Pt(2)
        
        # Title in primary color below the line (same as Introduction)
        title_top = Inches(6.8)  # Move to bottom
        title_textbox = slide.shapes.add_textbox(
            self.grid.get_position(0.5, 6.5, 1, 0.8)[0], title_top,
            Inches(6), Inches(0.8)
        )
        title_frame = title_textbox.text_frame
        title_frame.text = title
        self.styler.apply_title_style(title_frame.paragraphs[0], size=28, color=primary_color_hex)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    
    def _create_brand_story_slides(self, prs, identity_data, story_content):
        """Create Brand Story slides - split into multiple slides if content is too long"""
        
        # Split content into manageable chunks for slides
        story_parts = self._split_story_content(story_content)
        
        print(f"  📖 Brand Story split into {len(story_parts)} parts")
        for i, part in enumerate(story_parts):
            word_count = len(part.split())
            print(f"    Part {i+1}: {word_count} words")
        
        if len(story_parts) == 1:
            # Single slide
            self._create_brand_story_slide(prs, identity_data, "BRAND STORY", story_parts[0])
        else:
            # Multiple slides with pagination
            for i, part in enumerate(story_parts):
                slide_title = f"BRAND STORY ({i+1}/{len(story_parts)})"
                self._create_brand_story_slide(prs, identity_data, slide_title, part)
    
    def _split_story_content(self, content):
        """Split story content into appropriate chunks for slides"""
        if not content:
            return [""]
        
        # Target words per slide (roughly 150 words per slide for better readability on slides)
        words_per_slide = 150
        words = content.split()
        
        if len(words) <= words_per_slide:
            return [content]
        
        # Split into paragraphs first
        paragraphs = content.split('\n\n')
        
        parts = []
        current_part = ""
        current_word_count = 0
        
        for paragraph in paragraphs:
            paragraph_words = paragraph.split()
            paragraph_word_count = len(paragraph_words)
            
            # If adding this paragraph would exceed the limit, start a new part
            if current_word_count + paragraph_word_count > words_per_slide and current_part:
                parts.append(current_part.strip())
                current_part = paragraph + "\n\n"
                current_word_count = paragraph_word_count
            else:
                current_part += paragraph + "\n\n"
                current_word_count += paragraph_word_count
        
        # Add the last part if there's content
        if current_part.strip():
            parts.append(current_part.strip())
        
        # If we still don't have good splits, do a simple word-based split
        if not parts or (len(parts) == 1 and len(words) > words_per_slide * 2):
            parts = []
            for i in range(0, len(words), words_per_slide):
                chunk_words = words[i:i + words_per_slide]
                parts.append(' '.join(chunk_words))
        
        return parts if parts else [content]

    def create_pptx(self, company_name, identity_data, literature_data, brand_essence=None, 
                   industry="", values="", audience="", logo_color=None):
        """Create comprehensive brand book presentation with AI-generated content"""
        prs = Presentation()
        
        # Extract input parameters from brand_essence if not provided
        if brand_essence and brand_essence.get("company_profile"):
            profile = brand_essence["company_profile"]
            if not industry:
                industry = profile.get("industry", "technology")
            if not values and profile.get("core_values"):
                values = ", ".join(profile["core_values"])
            if not audience:
                audience = profile.get("target_audience", "businesses")
        
        # Research fonts
        print(f"🎨 Researching fonts for {company_name}...")
        self.researched_fonts = self.font_research_agent.research_fonts(company_name, industry, brand_essence)
        print(f"✅ Selected fonts: {self.researched_fonts['primary_font']} (primary), {self.researched_fonts['secondary_font']} (secondary)")
        
        # Generate enhanced color system BEFORE styling initialization
        print("🎨 Generating comprehensive color system with AI research...")
        try:
            color_system = self.enhanced_color_agent.generate_comprehensive_color_system(
                company_name, industry, values, audience, 
                brand_essence=brand_essence, logo_color=logo_color
            )
            self.enhanced_color_system = color_system
            print(f"✅ Enhanced color system generated with {len(color_system.get('primary_colors', []))} primary colors")
        except Exception as e:
            print(f"⚠️ Enhanced color generation failed, using fallback: {e}")
            color_system = self._generate_fallback_color_system(identity_data.get("palette", {}))
            self.enhanced_color_system = color_system
        
        # Initialize styling system AFTER color system is ready
        self._initialize_styling(identity_data, company_name)
        
        # Create slides
        sections = ["Table of Contents"]
        
        # 1. Title Slide
        self._create_title_slide(prs, company_name, identity_data, brand_essence)
        
        # Build sections list with Introduction and Brand Purpose sections
        sections.extend([
            "1. Introduction",
            "2. Brand Purpose (Vision)",
            "3. Brand Purpose (Mission)", 
            "4. Brand Purpose (Values)",
            "5. Brand Story",
            "6. Logo Variations",
            "7. Color Palette", 
            "8. Typography",
            "9. Imagery & Visuals",
        ])
        
        # 2. Table of Contents
        self._create_table_of_contents(prs, sections, identity_data)
        
        # 3. Introduction Slide with AI content
        self._create_introduction_slide(prs, company_name, brand_essence, identity_data, 
                                       industry, values, audience)
        
        # 4-6. Brand Purpose Slides (Vision, Mission, Core Values) with AI content
        self._create_brand_purpose_slides(prs, brand_essence, identity_data, 
                                         company_name, industry, values, audience)
        
        # 7. Brand Story Slide with AI content
        print("  🤖 Generating AI-powered brand story content...")
        try:
            story_content = self.brand_story_agent.generate_brand_story(
                company_name, industry, values, audience, brand_essence, literature_data
            )
        except Exception as e:
            print(f"  ⚠️ AI brand story generation failed, using fallback: {e}")
            story_content = literature_data.get("brand_story", "") if literature_data else ""
            if not story_content:
                story_content = f"Our story at {company_name} is one of innovation and dedication to excellence."
        
        if story_content:
            self._create_brand_story_slides(prs, identity_data, story_content)
        
        # 8. Logo Variations - RE-ENABLED AFTER FIXING ILLUSTRATION DISPLAY
        if identity_data.get("logos"):
            print(f"  🎨 Creating logo variations slide with {len(identity_data['logos'])} logos...")
            self._create_logo_variations_slide(prs, identity_data["logos"], identity_data)
        else:
            print("  ⚠️ No logos found in identity_data - skipping logo slide")
        
        # 9. Color Palette (3 slides: Primary, Secondary, Usage)
        if identity_data.get("palette"):
            self._create_color_palette_slide(prs, identity_data["palette"], identity_data, 
                                           company_name, industry, values, audience)
        
        # 10. Typography
        self._create_typography_slide(prs, identity_data.get("typography", {}), identity_data)
        
        # 9. Imagery & Visuals - Illustrations subsection  
        self._create_brand_illustrations_slide(prs, company_name, industry, values, audience, brand_essence, identity_data)
        
        # 9. Imagery & Visuals - Iconography subsection
        print("  🎨 Generating brand iconography with AI research and Fal AI...")
        try:
            primary_color_hex = self._get_primary_color_hex()
            
            # Generate Core Functional Icons using Fal.ai
            core_functional_categories = [
                "home icon", "search magnifying glass icon", "menu hamburger lines icon", 
                "settings gear icon", "user profile person icon", "notification bell icon", 
                "download arrow down icon", "upload arrow up icon"
            ]
            
            print("  🎨 Generating Core Functional Icons...")
            core_iconography = self.iconography_agent.create_iconography_system(
                company_name, industry, values, audience, primary_color_hex, 
                custom_categories=core_functional_categories, icon_type="core"
            )
            
            # Generate Industry-Specific Icons
            industry_categories = self._get_industry_icon_categories(industry)
            
            print(f"  🎨 Generating {industry.title()} Industry Icons...")
            industry_iconography = self.iconography_agent.create_iconography_system(
                company_name, industry, values, audience, primary_color_hex,
                custom_categories=industry_categories, icon_type="industry"
            )
            
            # Create slides with generated icons
            self._create_icons_display_slide(prs, core_iconography, identity_data, f"{company_name} - Core Functional Icons")
            self._create_icons_display_slide(prs, industry_iconography, identity_data, f"{company_name} - {industry.title()} Industry Icons")
            self._create_iconography_guidelines_slide(prs, core_iconography, identity_data)
            
            total_core = len(core_iconography.get('icon_generation', {}).get('generated_icons', []))
            total_industry = len(industry_iconography.get('icon_generation', {}).get('generated_icons', []))
            print(f"  ✅ Generated {total_core + total_industry} professional icons: {total_core} core + {total_industry} industry-specific")
        except Exception as e:
            print(f"  ⚠️ Iconography generation failed, creating basic slides: {e}")
            # Create basic iconography slides without generated icons
            self._create_icons_display_slide(prs, None, identity_data, company_name)
            self._create_iconography_guidelines_slide(prs, None, identity_data)
        
        # Brand apparel/merchandise slide - RE-ENABLED AFTER FIXING ILLUSTRATION DISPLAY
        if identity_data.get("logos") and len(identity_data["logos"]) > 0:
            first_logo_path = identity_data["logos"][0]  # Use first logo as reference
            print(f"  👕 Creating brand apparel slide using logo: {first_logo_path}")
            self._create_merchandise_slide_with_logo(prs, company_name, first_logo_path, identity_data)
        else:
            print("  ⚠️ No logos found - creating fallback brand apparel slide")
            self._create_merchandise_fallback_slide(prs, company_name, identity_data)
        
        # Brand mugs slide - RE-ENABLED AFTER FIXING ILLUSTRATION DISPLAY
        if identity_data.get("logos") and len(identity_data["logos"]) > 0:
            first_logo_path = identity_data["logos"][0]  # Use first logo as reference
            print(f"  ☕ Creating brand mugs slide using logo: {first_logo_path}")
            self._create_brand_mugs_slide(prs, company_name, first_logo_path, identity_data)
        else:
            print("  ⚠️ No logos found - creating fallback brand mugs slide")
            self._create_brand_mugs_fallback_slide(prs, company_name, identity_data)
        
        # Save file
        base_name = company_name.lower().replace(' ', '_')
        company_output_dir = os.path.join("output", base_name)
        os.makedirs(company_output_dir, exist_ok=True)
        file_name = f"{base_name}_enhanced_brand_book.pptx"
        file_path = os.path.join(company_output_dir, file_name)
        prs.save(file_path)
        
        print(f"Enhanced Brand Book PPTX saved at: {file_path}")
        return file_path


# For backward compatibility
PPTXGenerator = EnhancedPPTXGenerator

# Example usage
if __name__ == "__main__":
    identity_data = {
        "logos": [],
        "palette": {"primary": "#2E86AB", "secondary": "#A23B72", "accent": "#F18F01"},
        "typography": {"primary": "Inter", "secondary": "Source Sans Pro"},
        "visual_style": "Modern, clean, professional, minimalist, bold typography",
        "photography_style": "High contrast, authentic moments, natural lighting, diverse subjects"
    }
    
    literature_data = {
        "brand_story": "Our company revolutionizes the way people connect with technology through innovative design and user-centered solutions.",
        "voice_tone": "Professional yet approachable, confident, innovative, and human-centered.",
        "messaging_arch": "We believe technology should enhance human potential, not complicate it.",
        "marketing_copy": {
            "website": "Transform your digital experience with our cutting-edge solutions.",
            "social_media": "Innovation meets simplicity. #TechForHumans"
        },
        "collaterals": {
            "business_card": "Clean, modern design with QR code integration",
            "letterhead": "Professional header with brand elements"
        }
    }
    
    brand_essence = {
        "company_profile": {
            "name": "TechForward Inc",
            "industry": "Technology",
            "target_audience": "Progressive businesses seeking digital transformation",
            "core_values": ["Innovation", "Simplicity", "Human-Centered", "Excellence"]
        },
        "brand_positioning": {
            "unique_value_proposition": "We bridge the gap between complex technology and intuitive user experiences",
            "brand_promise": "Technology that works for you, not against you",
            "brand_personality": ["Innovative", "Reliable", "Approachable", "Progressive"],
            "competitive_advantage": "Our deep understanding of human psychology drives more intuitive technology solutions"
        }
    }
    
    generator = EnhancedPPTXGenerator()
    generator.create_pptx("TechForward Inc", identity_data, literature_data, brand_essence)
    
    
    
# done pptx.