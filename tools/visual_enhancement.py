try:
    from PIL import Image, ImageDraw, ImageFont, ImageFilter
except ImportError:
    Image = ImageDraw = ImageFont = ImageFilter = None

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.dml.color import RGBColor
    PPTX_AVAILABLE = True
except ImportError:
    Presentation = Inches = Pt = MSO_SHAPE = RGBColor = None
    PPTX_AVAILABLE = False

import colorsys
from typing import Dict, List, Tuple
import os

class VisualEnhancementEngine:
    """
    Advanced visual enhancement system for creating professional brand books
    with dark/gradient backgrounds, section dividers, and logo-derived shapes.
    """
    
    def __init__(self):
        self.gradient_styles = {
            "linear_diagonal": "diagonal",
            "radial_center": "radial",
            "linear_vertical": "vertical",
            "linear_horizontal": "horizontal"
        }
        
        self.section_divider_styles = {
            "geometric": "angular_shapes",
            "organic": "curved_shapes",
            "minimal": "clean_lines",
            "bold": "thick_bars"
        }
    
    def enhance_brand_book_visuals(self, palette: Dict, typography: Dict, 
                                  company_name: str, brand_style: str = "modern") -> Dict:
        """
        Create enhanced visual system with professional backgrounds and dividers
        """
        print(f"🎨 Enhancing visual system for {company_name}...")
        
        # Generate enhanced color system
        enhanced_palette = self._create_enhanced_palette(palette)
        
        # Create gradient backgrounds
        gradient_backgrounds = self._generate_gradient_backgrounds(enhanced_palette, brand_style)
        
        # Create section dividers
        section_dividers = self._create_section_dividers(enhanced_palette, brand_style)
        
        # Generate logo-derived shapes
        brand_shapes = self._create_brand_shapes(company_name, enhanced_palette)
        
        # Create visual hierarchy system
        visual_hierarchy = self._define_visual_hierarchy(enhanced_palette, typography)
        
        return {
            "enhanced_palette": enhanced_palette,
            "gradient_backgrounds": gradient_backgrounds,
            "section_dividers": section_dividers,
            "brand_shapes": brand_shapes,
            "visual_hierarchy": visual_hierarchy,
            "layout_grid": self._create_layout_grid(),
            "spacing_system": self._create_spacing_system()
        }
    
    def _create_enhanced_palette(self, base_palette: Dict) -> Dict:
        """Create extended palette with tints, shades, and gradients"""
        enhanced = {
            "base_colors": base_palette,
            "tints": {},
            "shades": {},
            "gradients": {},
            "accessible_pairs": []
        }
        
        # Extract base colors
        base_colors = {}
        if isinstance(base_palette, dict):
            for key, value in base_palette.items():
                if isinstance(value, str) and value.startswith('#'):
                    base_colors[key] = value
        
        # Generate tints and shades for each base color
        for name, hex_color in base_colors.items():
            enhanced["tints"][name] = self._generate_tints(hex_color)
            enhanced["shades"][name] = self._generate_shades(hex_color)
        
        # Create gradient combinations
        color_names = list(base_colors.keys())
        if len(color_names) >= 2:
            enhanced["gradients"] = {
                f"{color_names[0]}_to_{color_names[1]}": [base_colors[color_names[0]], base_colors[color_names[1]]],
                f"dark_to_{color_names[0]}": ["#1a1a1a", base_colors[color_names[0]]],
                f"{color_names[0]}_to_light": [base_colors[color_names[0]], "#f8f9fa"]
            }
        
        # Identify accessible color pairs
        enhanced["accessible_pairs"] = self._find_accessible_pairs(base_colors)
        
        return enhanced
    
    def _generate_tints(self, hex_color: str, steps: int = 5) -> List[str]:
        """Generate lighter tints of a color"""
        rgb = self._hex_to_rgb(hex_color)
        tints = []
        
        for i in range(1, steps + 1):
            factor = i / steps
            tinted_rgb = [
                int(rgb[j] + (255 - rgb[j]) * factor)
                for j in range(3)
            ]
            tints.append(self._rgb_to_hex(tinted_rgb))
        
        return tints
    
    def _generate_shades(self, hex_color: str, steps: int = 5) -> List[str]:
        """Generate darker shades of a color"""
        rgb = self._hex_to_rgb(hex_color)
        shades = []
        
        for i in range(1, steps + 1):
            factor = 1 - (i / steps)
            shaded_rgb = [int(rgb[j] * factor) for j in range(3)]
            shades.append(self._rgb_to_hex(shaded_rgb))
        
        return shades
    
    def _find_accessible_pairs(self, colors: Dict) -> List[Dict]:
        """Find accessible color combinations for text/background"""
        pairs = []
        color_list = list(colors.items())
        
        for i, (name1, color1) in enumerate(color_list):
            for j, (name2, color2) in enumerate(color_list):
                if i != j:
                    contrast = self._calculate_contrast_ratio(color1, color2)
                    if contrast >= 4.5:  # WCAG AA standard
                        pairs.append({
                            "foreground": name1,
                            "background": name2,
                            "contrast_ratio": contrast,
                            "wcag_aa": contrast >= 4.5,
                            "wcag_aaa": contrast >= 7.0
                        })
        
        return pairs[:10]  # Return top 10 pairs
    
    def _generate_gradient_backgrounds(self, palette: Dict, style: str) -> Dict:
        """Generate gradient background specifications"""
        base_colors = palette.get("base_colors", {})
        gradients = palette.get("gradients", {})
        
        background_gradients = {}
        
        # Hero section gradient
        if "primary" in base_colors and "accent" in base_colors:
            background_gradients["hero"] = {
                "type": "linear_diagonal",
                "colors": [base_colors["primary"], base_colors["accent"]],
                "direction": "135deg",
                "opacity": 0.9
            }
        
        # Section backgrounds
        background_gradients["section_light"] = {
            "type": "linear_vertical",
            "colors": ["#f8f9fa", "#ffffff"],
            "direction": "180deg",
            "opacity": 1.0
        }
        
        background_gradients["section_dark"] = {
            "type": "linear_vertical", 
            "colors": ["#1a1a1a", "#2d3748"],
            "direction": "180deg",
            "opacity": 1.0
        }
        
        # Accent gradients
        if "accent" in base_colors:
            accent_tints = palette.get("tints", {}).get("accent", [])
            if accent_tints:
                background_gradients["accent_subtle"] = {
                    "type": "radial_center",
                    "colors": [accent_tints[-1], "#ffffff"],
                    "direction": "center",
                    "opacity": 0.6
                }
        
        return background_gradients
    
    def _create_section_dividers(self, palette: Dict, style: str) -> Dict:
        """Create section divider specifications"""
        base_colors = palette.get("base_colors", {})
        
        dividers = {
            "geometric_bars": {
                "style": "rectangular",
                "height": "4px",
                "color": base_colors.get("accent", "#0066CC"),
                "margin": "2rem 0"
            },
            "gradient_divider": {
                "style": "linear_gradient",
                "height": "2px", 
                "gradient": f"linear-gradient(90deg, transparent, {base_colors.get('primary', '#333')}, transparent)",
                "margin": "3rem 0"
            },
            "brand_shape_divider": {
                "style": "custom_svg",
                "height": "20px",
                "pattern": "brand_derived",
                "color": base_colors.get("accent", "#0066CC")
            },
            "minimal_line": {
                "style": "solid",
                "height": "1px",
                "color": base_colors.get("secondary", "#666"),
                "opacity": 0.3,
                "margin": "1.5rem 0"
            }
        }
        
        return dividers
    
    def _create_brand_shapes(self, company_name: str, palette: Dict) -> Dict:
        """Generate logo-derived shapes and patterns"""
        
        # Extract first letter for monogram-style shapes
        first_letter = company_name[0].upper() if company_name else "B"
        
        shapes = {
            "monogram_circle": {
                "shape": "circle",
                "content": first_letter,
                "background": palette.get("base_colors", {}).get("primary", "#333"),
                "text_color": "#ffffff"
            },
            "geometric_pattern": {
                "shape": "hexagon_grid",
                "color": palette.get("base_colors", {}).get("accent", "#0066CC"),
                "opacity": 0.1,
                "usage": "background_pattern"
            },
            "brand_icon": {
                "shape": "rounded_rectangle",
                "aspect_ratio": "1:1",
                "corner_radius": "8px",
                "gradient": True
            },
            "decorative_elements": {
                "dots": {"color": palette.get("base_colors", {}).get("accent", "#0066CC"), "opacity": 0.6},
                "lines": {"color": palette.get("base_colors", {}).get("secondary", "#666"), "opacity": 0.3},
                "shapes": {"geometric": True, "organic": False}
            }
        }
        
        return shapes
    
    def _define_visual_hierarchy(self, palette: Dict, typography: Dict) -> Dict:
        """Define comprehensive visual hierarchy system"""
        base_colors = palette.get("base_colors", {})
        
        hierarchy = {
            "headings": {
                "h1": {
                    "font_size": "clamp(2.5rem, 5vw, 4rem)",
                    "font_weight": "bold",
                    "line_height": "1.1",
                    "color": base_colors.get("primary", "#333"),
                    "margin_bottom": "1.5rem"
                },
                "h2": {
                    "font_size": "clamp(2rem, 4vw, 3rem)",
                    "font_weight": "semibold", 
                    "line_height": "1.2",
                    "color": base_colors.get("primary", "#333"),
                    "margin_bottom": "1rem"
                },
                "h3": {
                    "font_size": "clamp(1.5rem, 3vw, 2rem)",
                    "font_weight": "medium",
                    "line_height": "1.3",
                    "color": base_colors.get("secondary", "#666"),
                    "margin_bottom": "0.75rem"
                }
            },
            "body_text": {
                "paragraph": {
                    "font_size": "1.125rem",
                    "line_height": "1.6",
                    "color": base_colors.get("secondary", "#666"),
                    "margin_bottom": "1rem"
                },
                "small_text": {
                    "font_size": "0.875rem",
                    "line_height": "1.5",
                    "color": "#888888",
                    "margin_bottom": "0.5rem"
                }
            },
            "emphasis": {
                "accent_text": {
                    "color": base_colors.get("accent", "#0066CC"),
                    "font_weight": "semibold"
                },
                "highlight": {
                    "background": palette.get("tints", {}).get("accent", ["#e6f3ff"])[0] if palette.get("tints") else "#e6f3ff",
                    "padding": "2px 6px",
                    "border_radius": "3px"
                }
            }
        }
        
        return hierarchy
    
    def _create_layout_grid(self) -> Dict:
        """Define responsive layout grid system"""
        return {
            "container": {
                "max_width": "1200px",
                "margin": "0 auto",
                "padding": "0 2rem"
            },
            "grid": {
                "columns": 12,
                "gap": "2rem",
                "breakpoints": {
                    "mobile": "768px",
                    "tablet": "1024px", 
                    "desktop": "1200px"
                }
            },
            "sections": {
                "hero": {"height": "100vh", "min_height": "600px"},
                "content": {"padding": "4rem 0"},
                "footer": {"padding": "2rem 0"}
            }
        }
    
    def _create_spacing_system(self) -> Dict:
        """Define consistent spacing system"""
        return {
            "scale": "1.25",  # Modular scale
            "base": "1rem",
            "sizes": {
                "xs": "0.25rem",
                "sm": "0.5rem", 
                "md": "1rem",
                "lg": "1.5rem",
                "xl": "2rem",
                "2xl": "3rem",
                "3xl": "4rem"
            },
            "components": {
                "button_padding": "0.75rem 1.5rem",
                "card_padding": "2rem",
                "section_margin": "4rem 0"
            }
        }
    
    def apply_visual_enhancements_to_pptx(self, prs: Presentation, 
                                         visual_system: Dict) -> Presentation:
        """Apply visual enhancements to PowerPoint presentation"""
        
        enhanced_palette = visual_system.get("enhanced_palette", {})
        gradients = visual_system.get("gradient_backgrounds", {})
        
        # Apply master slide styling
        for slide_layout in prs.slide_layouts:
            try:
                # Set background gradient if possible
                if hasattr(slide_layout, 'background'):
                    self._apply_gradient_to_slide_background(slide_layout, gradients.get("section_light"))
            except:
                pass  # Skip if background modification fails
        
        # Enhance existing slides
        for slide in prs.slides:
            self._enhance_slide_styling(slide, visual_system)
        
        return prs
    
    def _enhance_slide_styling(self, slide, visual_system: Dict):
        """Apply enhanced styling to individual slide"""
        enhanced_palette = visual_system.get("enhanced_palette", {})
        base_colors = enhanced_palette.get("base_colors", {})
        
        # Style title if present
        if hasattr(slide.shapes, 'title') and slide.shapes.title:
            title_shape = slide.shapes.title
            if hasattr(title_shape, 'text_frame'):
                for paragraph in title_shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.bold = True
                        if "primary" in base_colors:
                            primary_rgb = self._hex_to_rgb(base_colors["primary"])
                            run.font.color.rgb = RGBColor(*primary_rgb)
        
        # Add decorative elements
        self._add_decorative_elements_to_slide(slide, visual_system)
    
    def _add_decorative_elements_to_slide(self, slide, visual_system: Dict):
        """Add subtle decorative elements to slide"""
        try:
            enhanced_palette = visual_system.get("enhanced_palette", {})
            base_colors = enhanced_palette.get("base_colors", {})
            
            if "accent" in base_colors:
                # Add subtle accent line at top
                accent_rgb = self._hex_to_rgb(base_colors["accent"])
                line = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(0), Inches(0),
                    Inches(10), Inches(0.05)
                )
                line.fill.solid()
                line.fill.fore_color.rgb = RGBColor(*accent_rgb)
                line.line.fill.background()
                
        except Exception as e:
            pass  # Skip decorative elements if they fail
    
    def _hex_to_rgb(self, hex_color: str) -> Tuple[int, int, int]:
        """Convert hex color to RGB tuple"""
        hex_color = hex_color.lstrip('#')
        if len(hex_color) == 3:
            hex_color = ''.join([c*2 for c in hex_color])
        try:
            return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
        except:
            return (51, 51, 51)
    
    def _rgb_to_hex(self, rgb: List[int]) -> str:
        """Convert RGB to hex color"""
        return "#{:02x}{:02x}{:02x}".format(
            max(0, min(255, rgb[0])),
            max(0, min(255, rgb[1])), 
            max(0, min(255, rgb[2]))
        )
    
    def _calculate_contrast_ratio(self, color1: str, color2: str) -> float:
        """Calculate WCAG contrast ratio between two colors"""
        def luminance(hex_color: str) -> float:
            rgb = self._hex_to_rgb(hex_color)
            rgb_normalized = [c/255.0 for c in rgb]
            
            def adjust_gamma(c):
                return c / 12.92 if c <= 0.03928 else ((c + 0.055) / 1.055) ** 2.4
            
            rgb_linear = [adjust_gamma(c) for c in rgb_normalized]
            return 0.2126 * rgb_linear[0] + 0.7152 * rgb_linear[1] + 0.0722 * rgb_linear[2]
        
        lum1 = luminance(color1)
        lum2 = luminance(color2)
        
        lighter = max(lum1, lum2)
        darker = min(lum1, lum2)
        
        return (lighter + 0.05) / (darker + 0.05)