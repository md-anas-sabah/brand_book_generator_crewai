from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.text import MSO_AUTO_SIZE
import os
import re
import sys
sys.path.append(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
from agents.font_research_agent import FontResearchAgent

class PPTXGenerator:
    """
    Assembles brand assets and literature into a PowerPoint brand book.
    Now with intelligent font research integration.
    """
    
    def __init__(self):
        self.font_research_agent = FontResearchAgent()
        self.researched_fonts = None
    
    def _get_primary_font(self):
        """Get researched primary font or fallback"""
        return self.researched_fonts['primary_font'] if self.researched_fonts else "Inter"
    
    def _get_secondary_font(self):
        """Get researched secondary font or fallback"""
        return self.researched_fonts['secondary_font'] if self.researched_fonts else "Source Sans Pro"
    
    def _apply_font_to_paragraph(self, paragraph, font_type="primary", size=16, bold=False):
        """Apply researched font to any paragraph"""
        font_name = self._get_primary_font() if font_type == "primary" else self._get_secondary_font()
        paragraph.font.name = font_name
        paragraph.font.size = Pt(size)
        paragraph.font.bold = bold
    
    def _apply_font_to_textbox(self, textbox, font_type="primary", size=16, bold=False):
        """Apply researched font to all paragraphs in a textbox"""
        for paragraph in textbox.text_frame.paragraphs:
            self._apply_font_to_paragraph(paragraph, font_type, size, bold)
    
    def _style_title_text(self, paragraph, size=28, color_rgb=None):
        """Apply consistent title styling with researched font"""
        paragraph.font.name = self._get_primary_font()
        paragraph.font.size = Pt(size)
        paragraph.font.bold = True
        if color_rgb:
            paragraph.font.color.rgb = RGBColor(*color_rgb)
    
    def _style_content_text(self, paragraph, size=16, color_rgb=None):
        """Apply consistent content styling with researched font"""
        paragraph.font.name = self._get_secondary_font()
        paragraph.font.size = Pt(size)
        if color_rgb:
            paragraph.font.color.rgb = RGBColor(*color_rgb)

    def _add_title_slide(self, prs, company_name, identity_data=None, brand_essence=None):
        # Use blank layout for complete creative control
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # Add clean black background for modern minimalistic design
        self._add_black_background(slide)
        
        # Get brand colors for text
        palette = identity_data.get("palette", {}) if identity_data else {}
        text_color = self._get_text_color_for_black_bg(palette)
        text_rgb = self._hex_to_rgb(text_color)
        
        # Company name - centered, large and bold (modern minimalistic style)
        company_name_textbox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
        company_name_frame = company_name_textbox.text_frame
        company_name_frame.text = company_name
        self._style_title_text(company_name_frame.paragraphs[0], size=56, color_rgb=text_rgb)
        company_name_frame.paragraphs[0].alignment = 1  # Center alignment
        
        # Subtitle with modern styling - centered below company name
        subtitle_textbox = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(8), Inches(0.8))
        subtitle_frame = subtitle_textbox.text_frame
        subtitle_frame.text = "Brand Book"
        self._style_content_text(subtitle_frame.paragraphs[0], size=24, color_rgb=text_rgb)
        subtitle_frame.paragraphs[0].alignment = 1  # Center alignment
        
        # Add company logo if available - centered above company name
        self._add_title_logo(slide, identity_data)
    
    def _add_title_logo(self, slide, identity_data):
        """Add company logo to title slide in minimalistic style"""
        if not identity_data or not identity_data.get("logos"):
            return
            
        logos = identity_data["logos"]
        logo_path = None
        
        # Find the first existing logo
        for logo in logos:
            if isinstance(logo, str) and os.path.exists(logo):
                logo_path = logo
                break
        
        if logo_path:
            try:
                # Add logo centered above company name
                slide.shapes.add_picture(
                    logo_path, 
                    Inches(4), Inches(1),  # Centered horizontally, top area
                    width=Inches(2), height=Inches(1.2)
                )
            except Exception as e:
                print(f"Could not add logo to title slide: {e}")

    def _add_index_slide(self, prs, sections, identity_data=None, company_name=""):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Use blank layout
        
        # Add black background
        self._add_black_background(slide)
        
        # Get text color for black background (brand primary color)
        palette = identity_data.get("palette", {}) if identity_data else {}
        text_color = self._get_text_color_for_black_bg(palette)
        text_rgb = self._hex_to_rgb(text_color)
        
        # Add "Table of Contents" title
        title_textbox = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
        title_frame = title_textbox.text_frame
        title_frame.text = "Table of Contents"
        self._style_title_text(title_frame.paragraphs[0], size=32, color_rgb=text_rgb)
        title_frame.paragraphs[0].alignment = 1  # Center alignment
        
        # Create table of contents list (excluding "Table of Contents" from the list)
        content_lines = []
        for i, section in enumerate(sections[1:], 1):  # Skip first item which is "Table of Contents"
            content_lines.append(f"{i}. {section}")
        
        # Add content list
        content_textbox = slide.shapes.add_textbox(Inches(1.5), Inches(2), Inches(7), Inches(4.5))
        content_frame = content_textbox.text_frame
        content_frame.text = "\n".join(content_lines)
        content_frame.word_wrap = True
        
        # Style the content
        for paragraph in content_frame.paragraphs:
            self._style_content_text(paragraph, size=16, color_rgb=text_rgb)
    
    def _add_simple_index_slide(self, prs, identity_data=None, company_name=""):
        """Add the simple index slide with just company name and logo as requested"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Use blank layout
        
        # Add black background
        self._add_black_background(slide)
        
        # Get text color for black background (brand primary color)
        palette = identity_data.get("palette", {}) if identity_data else {}
        text_color = self._get_text_color_for_black_bg(palette)
        text_rgb = self._hex_to_rgb(text_color)
        
        # Add company name centered at top
        company_name_textbox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
        company_name_frame = company_name_textbox.text_frame
        company_name_frame.text = company_name
        company_name_frame.paragraphs[0].font.size = Pt(48)
        company_name_frame.paragraphs[0].font.bold = True
        company_name_frame.paragraphs[0].font.color.rgb = RGBColor(*text_rgb)
        company_name_frame.paragraphs[0].alignment = 1  # Center alignment
        
        # Add company logo if available - centered below company name
        if identity_data and identity_data.get("logos"):
            logos = identity_data["logos"]
            logo_path = None
            
            # Find the first existing logo
            for logo in logos:
                if isinstance(logo, str) and os.path.exists(logo):
                    logo_path = logo
                    break
            
            if logo_path:
                try:
                    # Add logo centered below company name
                    slide.shapes.add_picture(
                        logo_path, 
                        Inches(4), Inches(4),  # Centered horizontally
                        width=Inches(2)
                    )
                except Exception as e:
                    print(f"Could not add logo to index slide: {e}")

    def _get_brand_color(self, palette, color_name, default):
        """Extract color from palette safely"""
        if not palette:
            return default
        color = palette.get(color_name, default)
        if isinstance(color, list):
            return color[0] if color else default
        return color if color else default
    
    def _get_text_color_for_black_bg(self, palette):
        """Get appropriate text color for black background - use brand primary color"""
        # For black background, use the primary brand color for text
        primary_color = self._get_brand_color(palette, "primary", "#FFFFFF")
        
        try:
            # Calculate brightness to ensure it's not too dark on black
            hex_color = primary_color.lstrip('#')
            r = int(hex_color[0:2], 16)
            g = int(hex_color[2:4], 16) 
            b = int(hex_color[4:6], 16)
            brightness = (r * 299 + g * 587 + b * 114) / 1000
            
            # If primary color is too dark, use white instead
            if brightness < 80:  # Very dark colors
                return "#FFFFFF"
            else:
                return primary_color  # Use brand primary color
                
        except (ValueError, IndexError):
            # Fallback to white text
            return "#FFFFFF"
    
    def _hex_to_rgb(self, hex_color):
        """Convert hex color to RGB tuple"""
        hex_color = hex_color.lstrip('#')
        if len(hex_color) == 3:
            hex_color = ''.join([c*2 for c in hex_color])
        try:
            return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
        except:
            return (46, 134, 171)  # Default blue
    
    def _apply_text_color_to_textbox(self, textbox, color_hex):
        """Apply text color to all paragraphs and runs in a textbox - AGGRESSIVELY"""
        color_rgb = self._hex_to_rgb(color_hex)
        
        # Force color on ALL paragraphs
        for paragraph in textbox.text_frame.paragraphs:
            # Set paragraph-level color first
            paragraph.font.color.rgb = RGBColor(*color_rgb)
            
            # Then set run-level color for existing runs
            for run in paragraph.runs:
                run.font.color.rgb = RGBColor(*color_rgb)
            
            # If paragraph has no runs but has text, create a run
            if not paragraph.runs and paragraph.text.strip():
                paragraph.text = paragraph.text  # This creates a run
                paragraph.runs[0].font.color.rgb = RGBColor(*color_rgb)
        
        # Also set text_frame level defaults if available
        try:
            textbox.text_frame.paragraphs[0].font.color.rgb = RGBColor(*color_rgb)
        except:
            pass
    
    def _paginate_text(self, text, words_per_slide=50):
        """Split text into chunks for pagination"""
        if not text or not text.strip():
            return [""]
            
        words = text.split()
        if len(words) <= words_per_slide:
            return [text]
            
        chunks = []
        current_words = []
        
        for word in words:
            current_words.append(word)
            
            # Check if we've reached the word limit
            if len(current_words) >= words_per_slide:
                # Look for a natural break point (end of sentence)
                if word.endswith('.') or word.endswith('!') or word.endswith('?'):
                    chunks.append(' '.join(current_words))
                    current_words = []
                elif len(current_words) >= words_per_slide + 10:  # Force break if too long
                    chunks.append(' '.join(current_words))
                    current_words = []
        
        # Add remaining words if any
        if current_words:
            chunks.append(' '.join(current_words))
            
        return chunks
    
    def _create_gradient_background(self, slide, primary_color, accent_color):
        """Create a modern gradient background"""
        # Create a rectangle that covers the entire slide
        bg_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0, 0, 
            Inches(10), Inches(7.5)
        )
        
        # Set gradient fill
        fill = bg_shape.fill
        fill.gradient()
        fill.gradient_angle = 45  # Diagonal gradient
        
        # Set gradient stops
        gradient = fill.gradient_stops
        gradient[0].color.rgb = RGBColor(*self._hex_to_rgb(primary_color))
        gradient[1].color.rgb = RGBColor(*self._hex_to_rgb(accent_color))
        
        # Remove border
        bg_shape.line.fill.background()
    
    def _add_black_background(self, slide):
        """Add consistent black background to any slide"""
        # Create a rectangle that covers the entire slide
        bg_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0, 0, 
            Inches(10), Inches(7.5)
        )
        
        # Set solid black fill
        fill = bg_shape.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0, 0, 0)  # Pure black
        
        # Remove border
        bg_shape.line.fill.background()
        
        # Move background to back
        bg_shape.element.getparent().remove(bg_shape.element)
        slide.shapes._spTree.insert(1, bg_shape.element)
    
    def _add_geometric_elements(self, slide, accent_color, secondary_color):
        """Add modern geometric design elements"""
        # Large circle - top right
        circle1 = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(7), Inches(0.5),
            Inches(2.5), Inches(2.5)
        )
        fill = circle1.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*self._hex_to_rgb(secondary_color))
        circle1.fill.transparency = 0.3
        circle1.line.fill.background()
        
        # Small circle - bottom left  
        circle2 = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(0.5), Inches(5.5),
            Inches(1.5), Inches(1.5)
        )
        fill = circle2.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*self._hex_to_rgb(accent_color))
        circle2.fill.transparency = 0.4
        circle2.line.fill.background()
        
        # Rectangle accent
        rect = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(8.5), Inches(6),
            Inches(1), Inches(0.3)
        )
        fill = rect.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)
        rect.fill.transparency = 0.2
        rect.line.fill.background()
    
    def _add_hero_text(self, slide, text, primary="#FFFFFF", size=48, top=Inches(2), 
                      left=Inches(1), width=Inches(8), opacity=1.0):
        """Add styled hero text"""
        textbox = slide.shapes.add_textbox(left, top, width, Inches(1))
        text_frame = textbox.text_frame
        text_frame.text = text
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        
        # Style the paragraph
        paragraph = text_frame.paragraphs[0]
        paragraph.font.name = self._get_primary_font()
        paragraph.font.size = Pt(size)
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(*self._hex_to_rgb(primary))
        
        # Add transparency if needed
        if opacity < 1.0:
            textbox.fill.solid()
            textbox.fill.fore_color.rgb = RGBColor(*self._hex_to_rgb(primary))
            textbox.fill.transparency = 1.0 - opacity
    
    def _extract_brand_keywords(self, brand_essence, identity_data):
        """Extract key brand words from essence and identity"""
        keywords = []
        
        if brand_essence:
            # From brand positioning
            if brand_essence.get("brand_positioning", {}).get("brand_personality"):
                keywords.extend(brand_essence["brand_positioning"]["brand_personality"][:3])
            
            # From company profile values
            if brand_essence.get("company_profile", {}).get("core_values"):
                keywords.extend(brand_essence["company_profile"]["core_values"][:2])
        
        # From visual style if available
        if identity_data and identity_data.get("visual_style"):
            style_words = identity_data["visual_style"].split(",")[:2]
            keywords.extend([word.strip().title() for word in style_words])
        
        # Default keywords if none found
        if not keywords:
            keywords = ["Professional", "Innovative", "Trusted", "Modern"]
        
        return keywords[:4]  # Limit to 4 keywords
    
    def _add_brand_keywords(self, slide, keywords, accent_color):
        """Add brand keywords in a stylish layout"""
        if not keywords:
            return
            
        # Position keywords in a modern layout
        positions = [
            (Inches(1), Inches(4.5)),
            (Inches(3), Inches(4.8)), 
            (Inches(5), Inches(4.2)),
            (Inches(7), Inches(4.6))
        ]
        
        for i, keyword in enumerate(keywords[:4]):
            if i >= len(positions):
                break
                
            left, top = positions[i]
            
            # Create keyword box
            textbox = slide.shapes.add_textbox(left, top, Inches(1.8), Inches(0.5))
            text_frame = textbox.text_frame
            text_frame.text = keyword
            text_frame.word_wrap = False
            
            # Style the text
            paragraph = text_frame.paragraphs[0]
            paragraph.font.name = self._get_primary_font()
            paragraph.font.size = Pt(14)
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(255, 255, 255)
            
            # Add background shape
            bg_shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left - Inches(0.1), top - Inches(0.05),
                Inches(2), Inches(0.6)
            )
            fill = bg_shape.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(*self._hex_to_rgb(accent_color))
            bg_shape.fill.transparency = 0.8
            bg_shape.line.fill.background()
            
            # Send background behind text
            bg_shape.element.getparent().remove(bg_shape.element)
            slide.shapes._spTree.insert(2, bg_shape.element)
    
    def _add_hero_logo(self, slide, identity_data):
        """Add company logo to hero section if available"""
        if not identity_data or not identity_data.get("logos"):
            return
            
        logos = identity_data["logos"]
        logo_path = None
        
        # Find the first existing logo
        for logo in logos:
            if isinstance(logo, str) and os.path.exists(logo):
                logo_path = logo
                break
        
        if logo_path:
            try:
                # Add logo in top right corner
                slide.shapes.add_picture(
                    logo_path, 
                    Inches(7.5), Inches(1),
                    width=Inches(1.5)
                )
            except Exception as e:
                print(f"Could not add logo to title slide: {e}")
    
    def _add_bottom_decoration(self, slide, secondary_color):
        """Add decorative element at bottom of slide"""
        # Add a subtle bottom accent line
        line_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1), Inches(6.8),
            Inches(8), Inches(0.1)
        )
        fill = line_shape.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)
        line_shape.fill.transparency = 0.3
        line_shape.line.fill.background()

    def _add_logo_slide(self, prs, logos, identity_data=None):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        self._add_black_background(slide)
        
        # Get brand text color for black background
        palette = identity_data.get("palette", {}) if identity_data else {}
        text_color = self._get_text_color_for_black_bg(palette)
        text_rgb = self._hex_to_rgb(text_color)
        
        # Add title textbox
        title_textbox = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
        title_frame = title_textbox.text_frame
        title_frame.text = f"Logo Variations ({len(logos)} designs)"
        self._style_title_text(title_frame.paragraphs[0], size=28, color_rgb=text_rgb)
        
        left = Inches(1)
        top = Inches(2.5)
        width = Inches(2.2)
        height = Inches(2.2)
        
        print(f"Processing {len(logos)} logo variations...")
        
        for i, logo_path in enumerate(logos):
            if os.path.exists(logo_path):
                print(f"  Adding logo {i+1} with black background: {os.path.basename(logo_path)}")
                
                # Add black background rectangle for each logo
                logo_bg = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    left - Inches(0.15), top - Inches(0.15),
                    width + Inches(0.3), height + Inches(0.3)
                )
                logo_bg.fill.solid()
                logo_bg.fill.fore_color.rgb = RGBColor(0, 0, 0)  # Pure black
                logo_bg.line.fill.background()
                
                # Add logo on top of black background
                try:
                    slide.shapes.add_picture(logo_path, left, top, width=width, height=height)
                    print(f"  ✅ Logo {i+1} added successfully with black background")
                except Exception as e:
                    print(f"  ❌ Error adding logo {i+1}: {e}")
                
                # Move to next position
                left += width + Inches(0.4)
            else:
                print(f"  ❌ Logo file not found: {logo_path}")

    def _add_palette_slide(self, prs, palette, identity_data=None):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        self._add_black_background(slide)
        
        # Get brand text color for black background
        brand_palette = identity_data.get("palette", {}) if identity_data else {}
        text_color = self._get_text_color_for_black_bg(brand_palette)
        text_rgb = self._hex_to_rgb(text_color)
        
        # Add title textbox
        title_textbox = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
        title_frame = title_textbox.text_frame
        title_frame.text = "Color Palette"
        self._style_title_text(title_frame.paragraphs[0], size=28, color_rgb=text_rgb)
        
        left = Inches(1)
        top = Inches(1.5)
        width = Inches(1.5)
        height = Inches(1)
        spacing = Inches(0.3)
        for i, (color_name, hexcode) in enumerate(palette.items()):
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                left + i * (width + spacing),
                top,
                width,
                height
            )
            # Safe hex parsing - handle both strings and lists
            if isinstance(hexcode, list):
                hexcode = hexcode[0] if hexcode else "#CCCCCC"
            hexcode_clean = (hexcode or "#CCCCCC").lstrip("#")
            if len(hexcode_clean) == 3:
                hexcode_clean = ''.join([c*2 for c in hexcode_clean])
            if len(hexcode_clean) != 6:
                hexcode_clean = "CCCCCC"
            try:
                rgb = tuple(int(hexcode_clean[j:j+2], 16) for j in (0, 2, 4))
            except Exception:
                rgb = (204, 204, 204)
            fill = shape.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(*rgb)
            tf = shape.text_frame
            tf.text = f"{color_name.capitalize()}\n#{hexcode_clean.upper()}"
            for p in tf.paragraphs:
                p.font.size = Pt(12)
                p.font.bold = True
                # Use white text for better contrast on color swatches
                p.font.color.rgb = RGBColor(255, 255, 255)

    def _add_typography_slide(self, prs, typography, identity_data=None):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        self._add_black_background(slide)
        
        # Get brand text color for black background
        palette = identity_data.get("palette", {}) if identity_data else {}
        text_color = self._get_text_color_for_black_bg(palette)
        text_rgb = self._hex_to_rgb(text_color)
        
        # Add title textbox
        title_textbox = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
        title_frame = title_textbox.text_frame
        title_frame.text = "Typography"
        self._style_title_text(title_frame.paragraphs[0], size=28, color_rgb=text_rgb)
        
        left = Inches(1)
        top = Inches(1.5)
        width = Inches(6)
        height = Inches(2)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        tf.word_wrap = True
        # Show the researched fonts instead of the original typography data
        primary_font = self._get_primary_font()
        secondary_font = self._get_secondary_font()
        tf.text = f"Primary Font: {primary_font}\nSecondary Font: {secondary_font}"
        for p in tf.paragraphs:
            self._style_content_text(p, size=18, color_rgb=text_rgb)

    def _add_brand_essence_slides(self, prs, brand_essence, identity_data=None):
        """Add slides for brand essence and market analysis"""
        
        # Get brand text color for black background
        palette = identity_data.get("palette", {}) if identity_data else {}
        text_color = self._get_text_color_for_black_bg(palette)
        text_rgb = self._hex_to_rgb(text_color)
        
        # Company Profile Slide
        if brand_essence.get("company_profile"):
            profile = brand_essence["company_profile"]
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
            self._add_black_background(slide)
            
            # Add title textbox
            title_textbox = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
            title_frame = title_textbox.text_frame
            title_frame.text = "Company Profile"
            title_frame.paragraphs[0].font.size = Pt(28)
            title_frame.paragraphs[0].font.bold = True
            title_frame.paragraphs[0].font.color.rgb = RGBColor(*text_rgb)
            content = f"""
Company: {profile.get('name', 'N/A')}
Industry: {profile.get('industry', 'N/A')}
Target Audience: {profile.get('target_audience', 'N/A')}

Core Values:
{chr(10).join(['• ' + value for value in profile.get('core_values', [])])}
            """.strip()
            # Add content textbox
            content_textbox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(5))
            content_frame = content_textbox.text_frame
            content_frame.text = content
            content_frame.word_wrap = True
            # Apply brand text color to all content
            self._apply_text_color_to_textbox(content_textbox, text_color)
            text_rgb = self._hex_to_rgb(text_color)
            for paragraph in content_frame.paragraphs:
                paragraph.font.size = Pt(14)
                paragraph.font.color.rgb = RGBColor(*text_rgb)
        
        # Market Analysis Slide
        if brand_essence.get("market_analysis"):
            analysis = brand_essence["market_analysis"]
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
            self._add_black_background(slide)
            
            # Add title textbox
            title_textbox = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
            title_frame = title_textbox.text_frame
            title_frame.text = "Market Analysis & Insights"
            title_frame.paragraphs[0].font.size = Pt(28)
            title_frame.paragraphs[0].font.bold = True
            title_frame.paragraphs[0].font.color.rgb = RGBColor(*text_rgb)
            
            content_parts = []
            
            if analysis.get("industry_trends"):
                content_parts.append("Industry Trends:")
                content_parts.extend(['• ' + trend for trend in analysis["industry_trends"][:5]])
                content_parts.append("")
            
            if analysis.get("competitor_insights", {}).get("notable_competitors"):
                content_parts.append("Key Competitors:")
                competitors = analysis["competitor_insights"]["notable_competitors"][:6]
                content_parts.append(', '.join(competitors))
                content_parts.append("")
            
            if analysis.get("design_trends", {}).get("design_styles"):
                styles = analysis["design_trends"]["design_styles"][:4]
                content_parts.append(f"Popular Design Styles: {', '.join(styles)}")
            
            # Add content textbox
            content_textbox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(5))
            content_frame = content_textbox.text_frame
            content_frame.text = '\n'.join(content_parts)
            content_frame.word_wrap = True
            # Apply brand text color to all content
            self._apply_text_color_to_textbox(content_textbox, text_color)
            text_rgb = self._hex_to_rgb(text_color)
            for paragraph in content_frame.paragraphs:
                paragraph.font.size = Pt(14)
                paragraph.font.color.rgb = RGBColor(*text_rgb)
        
        # Brand Positioning Slide
        if brand_essence.get("brand_positioning"):
            positioning = brand_essence["brand_positioning"]
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
            self._add_black_background(slide)
            
            # Add title textbox
            title_textbox = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
            title_frame = title_textbox.text_frame
            title_frame.text = "Brand Positioning"
            title_frame.paragraphs[0].font.size = Pt(28)
            title_frame.paragraphs[0].font.bold = True
            title_frame.paragraphs[0].font.color.rgb = RGBColor(*text_rgb)
            
            content = f"""
Unique Value Proposition:
{positioning.get('unique_value_proposition', 'N/A')}

Brand Promise:
{positioning.get('brand_promise', 'N/A')}

Brand Personality:
{', '.join(positioning.get('brand_personality', []))}

Competitive Advantage:
{positioning.get('competitive_advantage', 'N/A')}
            """.strip()
            # Add content textbox
            content_textbox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(5))
            content_frame = content_textbox.text_frame
            content_frame.text = content
            content_frame.word_wrap = True
            # Apply brand text color to all content
            self._apply_text_color_to_textbox(content_textbox, text_color)
            text_rgb = self._hex_to_rgb(text_color)
            for paragraph in content_frame.paragraphs:
                paragraph.font.size = Pt(14)
                paragraph.font.color.rgb = RGBColor(*text_rgb)

    def _add_visual_style_slide(self, prs, visual_style, photography_style, identity_data=None):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        self._add_black_background(slide)
        
        # Get brand text color for black background
        palette = identity_data.get("palette", {}) if identity_data else {}
        text_color = self._get_text_color_for_black_bg(palette)
        text_rgb = self._hex_to_rgb(text_color)
        
        # Add title textbox
        title_textbox = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
        title_frame = title_textbox.text_frame
        title_frame.text = "Visual & Photography Guidelines"
        self._style_title_text(title_frame.paragraphs[0], size=28, color_rgb=text_rgb)
        
        left = Inches(1)
        top = Inches(1.5)
        width = Inches(7)
        height = Inches(2.5)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        tf.word_wrap = True
        tf.text = f"Visual Style:\n{visual_style}\n\nPhotography Style:\n{photography_style}"
        for p in tf.paragraphs:
            self._style_content_text(p, size=16, color_rgb=text_rgb)

    def _add_story_and_mission(self, prs, story, identity_data=None):
        # Attempt to split Brand Story, Mission, Values if text is present
        # Accepts either a dict (recommended) or a single string
        if isinstance(story, dict):
            if "Brand Story" in story:
                self._add_multislide_section(prs, "Brand Story", story["Brand Story"], 700, identity_data)
            if "Mission Statement" in story:
                self._add_multislide_section(prs, "Mission Statement", story["Mission Statement"], 700, identity_data)
            if "Our Values" in story:
                self._add_bullet_slide(prs, "Our Values", story["Our Values"], identity_data)
        else:
            # Fallback: Single slide, split if too long
            self._add_multislide_section(prs, "Brand Story & Mission", story, 900, identity_data)

    def _add_bullet_slide(self, prs, title, content, identity_data=None):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        self._add_black_background(slide)
        
        # Get brand text color for black background
        palette = identity_data.get("palette", {}) if identity_data else {}
        text_color = self._get_text_color_for_black_bg(palette)
        text_rgb = self._hex_to_rgb(text_color)
        
        # Add title textbox
        title_textbox = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
        title_frame = title_textbox.text_frame
        title_frame.text = title
        self._style_title_text(title_frame.paragraphs[0], size=28, color_rgb=text_rgb)
        
        left, top, width, height = Inches(1), Inches(1.5), Inches(7), Inches(3)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        tf.word_wrap = True
        # Accepts either a list or a string with \n, commas, or numbered
        if isinstance(content, list):
            for item in content:
                p = tf.add_paragraph()
                p.text = item.strip()
                p.level = 0
                self._style_content_text(p, size=18)
            # Apply contrasting text color to all content AGGRESSIVELY
            self._apply_text_color_to_textbox(textbox, text_color)
            text_rgb = self._hex_to_rgb(text_color)
            # Force color on every paragraph and run
            for p in tf.paragraphs:
                p.font.color.rgb = RGBColor(*text_rgb)
                for run in p.runs:
                    run.font.color.rgb = RGBColor(*text_rgb)
        else:
            # Try to split numbered/bulleted/commas
            items = []
            lines = content.split("\n")
            for line in lines:
                if not line.strip():
                    continue
                # Numbered?
                m = re.match(r"^\d+\.\s*(.*)", line.strip())
                if m:
                    items.append(m.group(1))
                elif "," in line and len(line.split(",")) <= 8:
                    items.extend([i.strip() for i in line.split(",")])
                else:
                    items.append(line.strip())
            for item in items:
                p = tf.add_paragraph()
                p.text = item.strip()
                p.level = 0
                self._style_content_text(p, size=18)
            # Apply contrasting text color to all content AGGRESSIVELY
            self._apply_text_color_to_textbox(textbox, text_color)
            text_rgb = self._hex_to_rgb(text_color)
            # Force color on every paragraph and run
            for p in tf.paragraphs:
                p.font.color.rgb = RGBColor(*text_rgb)
                for run in p.runs:
                    run.font.color.rgb = RGBColor(*text_rgb)

    def _add_multislide_section(self, prs, title, content, max_chars=800, identity_data=None):
        """Splits long text content into multiple slides if needed."""
        
        # Get contrasting text color for visibility
        palette = identity_data.get("palette", {}) if identity_data else {}
        text_color = self._get_text_color_for_black_bg(palette)
        text_rgb = self._hex_to_rgb(text_color)
        
        # Use centralized pagination logic
        chunks = self._paginate_text(content, 50)
        for idx, chunk in enumerate(chunks):
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
            self._add_black_background(slide)
            
            # Add title textbox
            title_textbox = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
            title_frame = title_textbox.text_frame
            title_frame.text = f"{title} ({idx+1})" if len(chunks) > 1 else title
            title_frame.paragraphs[0].font.size = Pt(28)
            title_frame.paragraphs[0].font.bold = True
            title_frame.paragraphs[0].font.color.rgb = RGBColor(*text_rgb)
            left, top, width, height = Inches(1), Inches(1.5), Inches(7), Inches(3)
            textbox = slide.shapes.add_textbox(left, top, width, height)
            tf = textbox.text_frame
            tf.word_wrap = True
            tf.text = chunk
            # Apply contrasting text color to all content AGGRESSIVELY
            self._apply_text_color_to_textbox(textbox, text_color)
            text_rgb = self._hex_to_rgb(text_color)
            for p in tf.paragraphs:
                size = 18 if len(chunk) < 300 else 16
                self._style_content_text(p, size=size, color_rgb=text_rgb)
                # Force color at run level
                for run in p.runs:
                    run.font.color.rgb = RGBColor(*text_rgb)

    def _add_voice_slide(self, prs, voice_tone, identity_data=None):
        self._add_multislide_section(prs, "Brand Voice & Tone", voice_tone, max_chars=750, identity_data=identity_data)

    def _add_messaging_slide(self, prs, messaging_arch, identity_data=None):
        # Try to split value propositions into bullets or slides with pagination
        lines = messaging_arch.split('\n')
        key_props = []
        for line in lines:
            m = re.match(r'^\d+\.\s*Key Value Proposition #[0-9]+: (.+)', line.strip())
            if m:
                key_props.append(m.group(1))
        if key_props:
            self._add_bullet_slide(prs, "Key Value Propositions", key_props, identity_data)
        else:
            self._add_multislide_section(prs, "Messaging & Value Propositions", messaging_arch, max_chars=700, identity_data=identity_data)

    def _add_marketing_copy_slides(self, prs, marketing_copy, identity_data=None):
        for channel, copy in marketing_copy.items():
            title = f"Marketing Copy: {channel.replace('_', ' ').title()}"
            # Use pagination with reasonable word count for marketing copy
            self._add_multislide_section(prs, title, copy, max_chars=600, identity_data=identity_data)

    def _add_collateral_slide(self, prs, collaterals, identity_data=None):
        # If it's a dict, bullet each with proper pagination
        if isinstance(collaterals, dict):
            for name, desc in collaterals.items():
                self._add_multislide_section(prs, f"Collateral: {name.replace('_',' ').title()}", desc, max_chars=600, identity_data=identity_data)
        else:
            self._add_multislide_section(prs, "Brand Collateral Templates", collaterals, max_chars=700, identity_data=identity_data)

    def create_pptx(self, company_name, identity_data, literature_data, brand_essence=None):
        prs = Presentation()
        
        # Research optimal fonts for this company
        print(f"🎨 Researching fonts for {company_name}...")
        industry = brand_essence.get("company_profile", {}).get("industry", "technology") if brand_essence else "technology"
        self.researched_fonts = self.font_research_agent.research_fonts(company_name, industry, brand_essence)
        print(f"✅ Selected fonts: {self.researched_fonts['primary_font']} (primary), {self.researched_fonts['secondary_font']} (secondary)")

        # Title Slide with modern design
        self._add_title_slide(prs, company_name, identity_data, brand_essence)
        
        # Create table of contents sections list
        sections = [
            "Table of Contents"
        ]
        
        if brand_essence:
            sections.extend([
                "Company Profile",
                "Market Analysis & Insights", 
                "Brand Positioning"
            ])
            
        sections.extend([
            "Logo Variations",
            "Color Palette", 
            "Typography",
            "Visual & Photography Guidelines",
            "Brand Story & Mission",
            "Brand Voice & Tone",
            "Messaging & Value Propositions",
            "Marketing Copy",
            "Brand Collateral Templates"
        ])
        
        # Add Table of Contents slide
        self._add_index_slide(prs, sections, identity_data, company_name)
        
        # Brand Essence & Market Analysis (if available)
        if brand_essence:
            self._add_brand_essence_slides(prs, brand_essence, identity_data)

        # Logo Variations
        self._add_logo_slide(prs, identity_data.get("logos", []), identity_data)

        # Color Palette
        self._add_palette_slide(prs, identity_data.get("palette", {}), identity_data)

        # Typography
        self._add_typography_slide(prs, identity_data.get("typography", {}), identity_data)

        # Visual Style & Photography
        self._add_visual_style_slide(
            prs,
            identity_data.get("visual_style", ""),
            identity_data.get("photography_style", ""),
            identity_data
        )

        # Brand Story, Mission, Values (handled as dict if possible)
        story_data = literature_data.get("brand_story", "")
        # Try to split brand_story content into dict
        if isinstance(story_data, str) and "Mission Statement:" in story_data and "Our Values:" in story_data:
            bd = {}
            m1 = re.search(r"Brand Story:(.+?)Mission Statement:", story_data, re.DOTALL)
            m2 = re.search(r"Mission Statement:(.+?)Our Values:", story_data, re.DOTALL)
            m3 = re.search(r"Our Values:(.+)", story_data, re.DOTALL)
            if m1: bd["Brand Story"] = m1.group(1).strip()
            if m2: bd["Mission Statement"] = m2.group(1).strip()
            if m3: bd["Our Values"] = m3.group(1).strip()
            self._add_story_and_mission(prs, bd, identity_data)
        else:
            self._add_story_and_mission(prs, story_data, identity_data)

        # Voice & Tone
        self._add_voice_slide(prs, literature_data.get("voice_tone", ""), identity_data)

        # Messaging
        self._add_messaging_slide(prs, literature_data.get("messaging_arch", ""), identity_data)

        # Marketing Copy
        self._add_marketing_copy_slides(prs, literature_data.get("marketing_copy", {}), identity_data)

        # Collateral
        self._add_collateral_slide(prs, literature_data.get("collaterals", {}), identity_data)

        # Save file in company-specific folder
        base_name = company_name.lower().replace(' ', '_')
        company_output_dir = os.path.join("output", base_name)
        os.makedirs(company_output_dir, exist_ok=True)
        file_name = f"{base_name}_brand_book.pptx"
        file_path = os.path.join(company_output_dir, file_name)
        prs.save(file_path)
        print(f"Brand Book PPTX saved at: {file_path}")
        return file_path

# Example usage
if __name__ == "__main__":
    identity_data = {
        "logos": [],
        "palette": {"primary": "#222", "secondary": "#EEE", "accent": "#FEC556"},
        "typography": {"primary": "Montserrat", "secondary": "Lato"},
        "visual_style": "Modern, minimal, lots of whitespace, bold accent color.",
        "photography_style": "Clean, authentic, natural light, focus on diversity."
    }
    literature_data = {
        "brand_story": {
            "Brand Story": "Acme is here to innovate fintech for everyone.",
            "Mission Statement": "Make finance easy for all.",
            "Our Values": ["Innovation", "Customer Focus", "Integrity"]
        },
        "voice_tone": "Friendly, helpful, smart.",
        "messaging_arch": "1. Key Value Proposition #1: Fast\n2. Key Value Proposition #2: Secure",
        "marketing_copy": {"website": "Welcome to Acme!", "email": "Acme for your finance.", "social_media": "Finance made easy."},
        "collaterals": {"business_card": "Modern layout", "email_signature": "Professional details"}
    }
    generator = PPTXGenerator()
    generator.create_pptx("Acme Corp", identity_data, literature_data)
