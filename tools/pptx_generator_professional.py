from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.enum.dml import MSO_THEME_COLOR
import os
import re

class ProfessionalPPTXGenerator:
    """
    Professional-grade PowerPoint generator matching top agency standards.
    Features clean minimalistic design, sophisticated typography, and brand consistency.
    """

    def __init__(self):
        # Professional color scheme - using sophisticated palette
        self.bg_color = (0, 0, 0)  # Pure black background
        self.primary_text = (255, 255, 255)  # Pure white text
        self.accent_color = None  # Will be set from brand colors
        self.secondary_text = (180, 180, 180)  # Light gray for secondary text
        
        # Typography settings
        self.title_font = "Montserrat"
        self.body_font = "Open Sans"
        self.accent_font = "Montserrat"

    def _get_brand_accent_color(self, palette):
        """Extract and set brand accent color"""
        if not palette:
            self.accent_color = (255, 235, 59)  # Default yellow
            return
        
        # Priority order for accent color selection
        color_priorities = ["accent", "primary", "secondary"]
        
        for color_name in color_priorities:
            color = palette.get(color_name)
            if color:
                if isinstance(color, list):
                    color = color[0] if color else None
                if color and color.startswith('#'):
                    self.accent_color = self._hex_to_rgb(color)
                    return
        
        # Fallback to sophisticated yellow
        self.accent_color = (255, 235, 59)

    def _hex_to_rgb(self, hex_color):
        """Convert hex to RGB tuple"""
        hex_color = hex_color.lstrip('#')
        if len(hex_color) == 3:
            hex_color = ''.join([c*2 for c in hex_color])
        try:
            return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
        except:
            return (255, 235, 59)  # Default yellow

    def _create_black_slide(self, prs):
        """Create slide with professional black background"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # Fill entire slide with black background
        bg_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0, 0, 
            Inches(10), Inches(7.5)
        )
        fill = bg_shape.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*self.bg_color)
        bg_shape.line.fill.background()
        
        return slide

    def _add_accent_line(self, slide, y_position=Inches(6.8)):
        """Add professional yellow accent line"""
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.8), y_position,
            Inches(8.4), Inches(0.05)
        )
        fill = line.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*self.accent_color)
        line.line.fill.background()

    def _add_page_number(self, slide, page_num):
        """Add professional page numbering"""
        page_box = slide.shapes.add_textbox(
            Inches(8.5), Inches(7), 
            Inches(1), Inches(0.3)
        )
        text_frame = page_box.text_frame
        text_frame.text = f"{page_num:02d}"
        
        paragraph = text_frame.paragraphs[0]
        paragraph.font.name = self.body_font
        paragraph.font.size = Pt(14)
        paragraph.font.color.rgb = RGBColor(*self.accent_color)
        paragraph.font.bold = True
        paragraph.alignment = PP_ALIGN.RIGHT

    def _add_section_title(self, slide, title, y_position=Inches(6.2)):
        """Add professional section title in accent color"""
        title_box = slide.shapes.add_textbox(
            Inches(0.8), y_position,
            Inches(8), Inches(0.8)
        )
        text_frame = title_box.text_frame
        text_frame.text = title.upper()
        
        paragraph = text_frame.paragraphs[0]
        paragraph.font.name = self.accent_font
        paragraph.font.size = Pt(36)
        paragraph.font.color.rgb = RGBColor(*self.accent_color)
        paragraph.font.bold = True
        paragraph.alignment = PP_ALIGN.LEFT

    def _add_main_content(self, slide, content, y_start=Inches(1.2)):
        """Add main content with professional typography"""
        content_box = slide.shapes.add_textbox(
            Inches(0.8), y_start,
            Inches(8), Inches(4.5)
        )
        text_frame = content_box.text_frame
        text_frame.text = content
        text_frame.word_wrap = True
        
        paragraph = text_frame.paragraphs[0]
        paragraph.font.name = self.body_font
        paragraph.font.size = Pt(24)
        paragraph.font.color.rgb = RGBColor(*self.primary_text)
        paragraph.line_spacing = 1.4
        paragraph.alignment = PP_ALIGN.LEFT

    def _create_title_slide(self, prs, company_name, identity_data=None):
        """Create sophisticated title slide matching professional standards"""
        slide = self._create_black_slide(prs)
        
        # Set brand colors
        palette = identity_data.get("palette", {}) if identity_data else {}
        self._get_brand_accent_color(palette)
        
        # Main company name - large, bold, centered
        title_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(2.5),
            Inches(8.4), Inches(2)
        )
        text_frame = title_box.text_frame
        text_frame.text = company_name.upper()
        
        paragraph = text_frame.paragraphs[0]
        paragraph.font.name = self.title_font
        paragraph.font.size = Pt(72)
        paragraph.font.color.rgb = RGBColor(*self.primary_text)
        paragraph.font.bold = True
        paragraph.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(4.8),
            Inches(8.4), Inches(0.8)
        )
        text_frame = subtitle_box.text_frame
        text_frame.text = "BRAND IDENTITY SYSTEM"
        
        paragraph = text_frame.paragraphs[0]
        paragraph.font.name = self.accent_font
        paragraph.font.size = Pt(24)
        paragraph.font.color.rgb = RGBColor(*self.accent_color)
        paragraph.font.bold = True
        paragraph.alignment = PP_ALIGN.CENTER
        
        # Add accent line
        self._add_accent_line(slide)
        
        return slide

    def _create_index_slide(self, prs, sections):
        """Create professional index page"""
        slide = self._create_black_slide(prs)
        
        # Add section title
        self._add_section_title(slide, "INDEX")
        
        # Create two-column layout for index
        left_sections = sections[::2]  # Even indices
        right_sections = sections[1::2]  # Odd indices
        
        # Left column
        left_content = []
        for i, section in enumerate(left_sections):
            page_num = (i * 2) + 1
            left_content.append(f"{section:<30} {page_num:02d}")
        
        left_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(1.5),
            Inches(4), Inches(4)
        )
        text_frame = left_box.text_frame
        text_frame.text = "\n".join(left_content)
        
        paragraph = text_frame.paragraphs[0]
        paragraph.font.name = self.body_font
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = RGBColor(*self.primary_text)
        paragraph.line_spacing = 1.6
        
        # Right column
        if right_sections:
            right_content = []
            for i, section in enumerate(right_sections):
                page_num = (i * 2) + 2
                right_content.append(f"{section:<30} {page_num:02d}")
            
            right_box = slide.shapes.add_textbox(
                Inches(5.2), Inches(1.5),
                Inches(4), Inches(4)
            )
            text_frame = right_box.text_frame
            text_frame.text = "\n".join(right_content)
            
            paragraph = text_frame.paragraphs[0]
            paragraph.font.name = self.body_font
            paragraph.font.size = Pt(18)
            paragraph.font.color.rgb = RGBColor(*self.primary_text)
            paragraph.line_spacing = 1.6
        
        # Add accent line and page number
        self._add_accent_line(slide)
        self._add_page_number(slide, 1)

    def _create_content_slide(self, prs, title, content, page_num):
        """Create professional content slide"""
        slide = self._create_black_slide(prs)
        
        # Add section title
        self._add_section_title(slide, title)
        
        # Add main content
        self._add_main_content(slide, content)
        
        # Add accent line and page number
        self._add_accent_line(slide)
        self._add_page_number(slide, page_num)
        
        return slide

    def _create_logo_slide(self, prs, logos, page_num):
        """Create professional logo showcase slide"""
        slide = self._create_black_slide(prs)
        
        # Add section title
        self._add_section_title(slide, "LOGO")
        
        # Add logos if available
        if logos:
            logo_added = False
            for i, logo_path in enumerate(logos[:3]):  # Max 3 logos
                if isinstance(logo_path, str) and os.path.exists(logo_path):
                    try:
                        x_pos = Inches(1.5 + (i * 2.5))
                        slide.shapes.add_picture(
                            logo_path, 
                            x_pos, Inches(2.5),
                            width=Inches(2)
                        )
                        logo_added = True
                    except Exception as e:
                        continue
            
            if not logo_added:
                # Add placeholder text
                self._add_main_content(slide, "Logo variations will be displayed here once generated.", Inches(2.5))
        else:
            self._add_main_content(slide, "Logo variations will be displayed here once generated.", Inches(2.5))
        
        # Add accent line and page number
        self._add_accent_line(slide)
        self._add_page_number(slide, page_num)

    def _create_color_palette_slide(self, prs, palette, page_num):
        """Create professional color palette slide"""
        slide = self._create_black_slide(prs)
        
        # Add section title
        self._add_section_title(slide, "BRAND COLORS")
        
        if palette:
            # Create color swatches
            x_start = Inches(1)
            y_start = Inches(2.5)
            swatch_width = Inches(1.8)
            swatch_height = Inches(1.8)
            spacing = Inches(0.3)
            
            colors_per_row = 4
            row = 0
            col = 0
            
            for name, color_value in list(palette.items())[:8]:  # Max 8 colors
                if isinstance(color_value, str) and color_value.startswith('#'):
                    x_pos = x_start + col * (swatch_width + spacing)
                    y_pos = y_start + row * (swatch_height + spacing + Inches(0.5))
                    
                    # Color swatch
                    swatch = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        x_pos, y_pos,
                        swatch_width, swatch_height
                    )
                    fill = swatch.fill
                    fill.solid()
                    fill.fore_color.rgb = RGBColor(*self._hex_to_rgb(color_value))
                    swatch.line.fill.background()
                    
                    # Color name and value
                    label_box = slide.shapes.add_textbox(
                        x_pos, y_pos + swatch_height + Inches(0.1),
                        swatch_width, Inches(0.6)
                    )
                    text_frame = label_box.text_frame
                    text_frame.text = f"{name.upper()}\n{color_value.upper()}"
                    
                    paragraph = text_frame.paragraphs[0]
                    paragraph.font.name = self.body_font
                    paragraph.font.size = Pt(12)
                    paragraph.font.color.rgb = RGBColor(*self.primary_text)
                    paragraph.font.bold = True
                    paragraph.alignment = PP_ALIGN.CENTER
                    
                    col += 1
                    if col >= colors_per_row:
                        col = 0
                        row += 1
        else:
            self._add_main_content(slide, "Brand color palette will be displayed here.", Inches(2.5))
        
        # Add accent line and page number
        self._add_accent_line(slide)
        self._add_page_number(slide, page_num)

    def create_professional_pptx(self, company_name, identity_data, literature_data, brand_essence=None):
        """Create professional-grade PowerPoint presentation"""
        prs = Presentation()
        page_count = 1
        
        # Set brand colors
        palette = identity_data.get("palette", {}) if identity_data else {}
        self._get_brand_accent_color(palette)
        
        # Title slide
        self._create_title_slide(prs, company_name, identity_data)
        
        # Define sections for index
        sections = []
        if brand_essence:
            sections.extend(["Introduction", "Brand Purpose", "Brand Story"])
        
        sections.extend([
            "Logo", "Brand Colors", "Typography", 
            "Visual Language", "Brand Voice", "Applications"
        ])
        
        # Index slide
        self._create_index_slide(prs, sections)
        page_count += 1
        
        # Brand essence slides
        if brand_essence:
            # Introduction
            if brand_essence.get("company_profile"):
                profile = brand_essence["company_profile"]
                intro_content = f"""The following brand identity system for {company_name} is thoughtfully crafted to present our brand in an international, engaging, consistent, recognisable and proprietary way. Unique in form, versatile in its application and unified by a fundamental principle.

Industry: {profile.get('industry', 'N/A')}
Target Audience: {profile.get('target_audience', 'N/A')}"""
                
                self._create_content_slide(prs, "INTRODUCTION", intro_content, page_count)
                page_count += 1
            
            # Brand Purpose
            if brand_essence.get("brand_positioning"):
                positioning = brand_essence["brand_positioning"]
                purpose_content = positioning.get("unique_value_proposition", 
                    "We are here to enable competitive excellence and empower aspiring professionals with our world-class platform.")
                
                self._create_content_slide(prs, "BRAND PURPOSE", purpose_content, page_count)
                page_count += 1
        
        # Logo slide
        self._create_logo_slide(prs, identity_data.get("logos", []), page_count)
        page_count += 1
        
        # Color palette slide
        self._create_color_palette_slide(prs, palette, page_count)
        page_count += 1
        
        # Typography slide
        if identity_data.get("typography"):
            typography = identity_data["typography"]
            typo_content = f"""Primary Font: {typography.get('primary', 'Montserrat')}
Secondary Font: {typography.get('secondary', 'Open Sans')}

{typography.get('description', 'Professional typography system designed for maximum readability and brand consistency across all applications.')}"""
            
            self._create_content_slide(prs, "TYPOGRAPHY", typo_content, page_count)
            page_count += 1
        
        # Brand voice slide
        if literature_data and literature_data.get("voice_tone"):
            voice_content = literature_data["voice_tone"]
            self._create_content_slide(prs, "BRAND VOICE", voice_content, page_count)
            page_count += 1
        
        # Save file
        base_name = company_name.lower().replace(' ', '_')
        company_output_dir = os.path.join("output", base_name)
        os.makedirs(company_output_dir, exist_ok=True)
        file_name = f"{base_name}_brand_book_professional.pptx"
        file_path = os.path.join(company_output_dir, file_name)
        prs.save(file_path)
        print(f"Professional Brand Book PPTX saved at: {file_path}")
        return file_path

# Replace the old generator
PPTXGenerator = ProfessionalPPTXGenerator